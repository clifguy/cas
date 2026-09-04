"""SAGE adapter tests (TEST-SAGE-AD-001 through AD-199).

Production adapter tests for nomic-embed-text EmbeddingProvider, Qwen3
AbstractionProvider (with lazy loading), and Markdown source adapter
provenance. Embedding tests require nomic-embed-text (~270MB download on
first run). Abstraction tests require mlx-lm and the weights of the model
``sage/config.yaml`` ships (several GB on first run); they form the opt-in
real-model tier, skipped unless ``SAGE_TEST_REAL_MODELS=1`` (see
``tests/helpers/real_models.py``), which keeps the default run fast and
light. Markdown adapter tests have no external dependencies. Content-store
adapter tests live in
tests/sage/test_content_store_postgres.py.

Tests are organized in implementation dependency order: embedding provider
first, then abstraction provider, then markdown adapter.
"""

import math
from collections.abc import Awaitable, Callable, Iterator
from pathlib import Path

import pytest

from sage.config import load_sage_core_config
from tests.helpers.real_models import (
    loaded_provider,
    real_model_lock,
    release_provider,
    requires_real_models,
)

# ── Skip if dependencies unavailable ────────────────────────────────

try:
    from sage.adapters.embedding_nomic import NomicEmbeddingProvider

    _HAS_EMBEDDING = True
except (ImportError, RuntimeError):
    _HAS_EMBEDDING = False

try:
    import mlx_lm  # noqa: F401 # gate fires on missing runtime dep, not just module shape

    from sage.adapters.abstraction_qwen3 import Qwen3AbstractionProvider

    _HAS_QWEN3 = True
except (ImportError, RuntimeError):
    _HAS_QWEN3 = False


requires_embedding = pytest.mark.skipif(
    not _HAS_EMBEDDING, reason="sentence-transformers or nomic model not available"
)
requires_qwen3 = pytest.mark.skipif(
    not _HAS_QWEN3,
    reason=(
        "Qwen3 abstraction tests require mlx-lm (Apple Silicon only); "
        "skipped on Linux CI runners by design"
    ),
)


# ── Fixtures ────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def embedding_provider():
    """Module-scoped embedding provider (model loads once for all tests)."""
    if not _HAS_EMBEDDING:
        pytest.skip("sentence-transformers or nomic model not available")
    return NomicEmbeddingProvider()


# ══════════════════════════════════════════════════════════════════════
# 1. nomic-embed-text EmbeddingProvider
# ══════════════════════════════════════════════════════════════════════


@requires_embedding
class TestNomicEmbeddingProvider:
    """Tests AD-001 through AD-008."""

    async def test_ad_001_embedding_dimension_is_768(self, embedding_provider):
        """AD-001: nomic-embed-text produces 768-dimensional vectors."""
        result = await embedding_provider.embed(
            ["The report claims a novel method for data synchronization."]
        )
        assert len(result) == 1
        assert len(result[0]) == 768
        assert all(isinstance(x, float) for x in result[0])

    async def test_ad_002_batch_preserves_order(self, embedding_provider):
        """AD-002: N input texts produce N output vectors in order."""
        texts = ["alpha", "beta", "gamma"]
        result = await embedding_provider.embed(texts)
        assert len(result) == 3
        for vec in result:
            assert len(vec) == 768

        # Verify order: re-embed individually and compare
        for i, text in enumerate(texts):
            single = await embedding_provider.embed([text])
            # Should match the corresponding batch vector
            diff = sum((a - b) ** 2 for a, b in zip(result[i], single[0]))
            assert diff < 1e-6, f"Batch vector {i} does not match single embed"

    async def test_ad_003_deterministic_embeddings(self, embedding_provider):
        """AD-003: Same text produces identical embeddings."""
        text = ["reproducibility test"]
        result1 = await embedding_provider.embed(text)
        result2 = await embedding_provider.embed(text)

        diff = sum((a - b) ** 2 for a, b in zip(result1[0], result2[0]))
        assert diff < 1e-6

    async def test_ad_004_l2_normalized(self, embedding_provider):
        """AD-004: Output vectors are L2-normalized (norm ~= 1.0)."""
        texts = [
            "Short text.",
            "A significantly longer passage with multiple sentences about "
            "various topics including report law, data management, and "
            "retrieval systems.",
        ]
        result = await embedding_provider.embed(texts)

        for vec in result:
            norm = math.sqrt(sum(x * x for x in vec))
            assert abs(norm - 1.0) < 1e-4, f"L2 norm = {norm}, expected ~1.0"

    async def test_ad_005_semantic_similarity(self, embedding_provider):
        """AD-005: Similar texts have higher cosine similarity than dissimilar texts."""
        texts = [
            "The document describes a method for synchronizing health records.",
            "A technique for keeping medical data in sync across systems.",
            "The basketball team scored 47 points in the first half.",
        ]
        vecs = await embedding_provider.embed(texts)

        sim_ab = _cosine_sim(vecs[0], vecs[1])
        sim_ac = _cosine_sim(vecs[0], vecs[2])
        sim_bc = _cosine_sim(vecs[1], vecs[2])

        assert sim_ab > sim_ac, f"sim(A,B)={sim_ab:.4f} should > sim(A,C)={sim_ac:.4f}"
        assert sim_ab > sim_bc, f"sim(A,B)={sim_ab:.4f} should > sim(B,C)={sim_bc:.4f}"

    async def test_ad_006_empty_input(self, embedding_provider):
        """AD-006: embed([]) returns [] immediately."""
        result = await embedding_provider.embed([])
        assert result == []

    async def test_ad_007_single_input(self, embedding_provider):
        """AD-007: A batch of one text works correctly."""
        result = await embedding_provider.embed(["single input"])
        assert len(result) == 1
        assert len(result[0]) == 768
        norm = math.sqrt(sum(x * x for x in result[0]))
        assert abs(norm - 1.0) < 1e-4

    def test_ad_008_init_fails_on_bad_model(self):
        """AD-008: Provider init fails fast if model unavailable."""
        with pytest.raises(RuntimeError, match="nonexistent-model-xyz"):
            NomicEmbeddingProvider(model_name="nonexistent-model-xyz")


async def test_stub_content_store_reports_zero_retained_versions():
    """StubContentStore has no on-disk versioning, so count_retained_versions is 0.

    The substrate-agnostic contract: callers routing through the ContentStore
    interface receive a well-formed integer without special-casing the stub.
    """
    from sage.adapters.stubs import StubContentStore

    store = StubContentStore()
    assert await store.count_retained_versions() == 0


async def test_stub_content_store_reports_zero_small_fragments():
    """StubContentStore has no on-disk fragments, so count_small_fragments is 0.

    Same substrate-agnostic contract as the retained-version measure.
    """
    from sage.adapters.stubs import StubContentStore

    store = StubContentStore()
    assert await store.count_small_fragments() == 0


# ── Helpers ─────────────────────────────────────────────────────────


def _cosine_sim(a: list[float], b: list[float]) -> float:
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = math.sqrt(sum(x * x for x in a))
    norm_b = math.sqrt(sum(x * x for x in b))
    if norm_a == 0.0 or norm_b == 0.0:
        return 0.0
    return dot / (norm_a * norm_b)


# ══════════════════════════════════════════════════════════════════════
# 3. Qwen3 AbstractionProvider
# ══════════════════════════════════════════════════════════════════════

_PROJECT_ROOT = Path(__file__).resolve().parents[2]


def committed_qwen3_model_id() -> str:
    """The abstraction model the stack config ships, through the validated loader.

    Read lazily so a config problem surfaces in the real-model tier rather
    than failing this whole module at collection, and checked to be an MLX
    model so a committed switch to a hosted provider cannot hand its model id
    to MLX.
    """
    config = load_sage_core_config(_PROJECT_ROOT / "sage" / "config.yaml")
    if config.abstraction.provider != "local-mlx":
        raise RuntimeError(
            f"stack abstraction provider is {config.abstraction.provider!r}, not "
            "'local-mlx'; the committed model id is not an MLX model"
        )
    if not config.abstraction.model:
        raise RuntimeError("stack config declares no abstraction.model")
    return config.abstraction.model


def _assert_released(provider: Qwen3AbstractionProvider) -> None:
    """Fail the class at teardown if its provider still holds model state."""
    assert provider._model is None, "provider still holds a model after release"
    assert provider._tokenizer is None, "provider still holds a tokenizer after release"
    assert provider._executor is None, "provider still holds its executor after release"


SAMPLE_TEXT = (
    "The document describes a method for synchronizing patient health records "
    "across distributed clinical systems using a conflict-free replicated data "
    "type (CRDT). The approach ensures eventual consistency while preserving "
    "the causal ordering of clinical events. The system handles concurrent "
    "updates from multiple hospital sites without requiring a central "
    "coordinator, reducing single points of failure in the health information "
    "exchange infrastructure."
)


@pytest.fixture(scope="class")
def qwen3_provider() -> Iterator[Qwen3AbstractionProvider]:
    """Class-scoped abstraction provider. With lazy loading, construction is
    cheap (no model allocated); the first generate_abstract() call loads the
    model. Held under the machine-wide real-model lock, and unloaded when the
    class finishes so the weights do not stay resident for the rest of the
    run; the teardown asserts the release actually happened."""
    if not _HAS_QWEN3:
        pytest.skip("mlx-lm or Qwen3 model not available")
    with loaded_provider(
        lambda: Qwen3AbstractionProvider(model_id=committed_qwen3_model_id())
    ) as provider:
        yield provider
    _assert_released(provider)


@pytest.fixture(scope="class")
def qwen3_provider_factory() -> Iterator[Callable[..., Awaitable[Qwen3AbstractionProvider]]]:
    """Builds providers on demand for the lazy-loading tests, under the same
    lock. Only one provider is ever loaded at a time: building a new one first
    unloads the earlier ones, and every provider built is unloaded again when
    the class finishes, which the teardown asserts."""
    if not _HAS_QWEN3:
        pytest.skip("mlx-lm or Qwen3 model not available")
    created: list[Qwen3AbstractionProvider] = []

    async def make(model_id: str | None = None) -> Qwen3AbstractionProvider:
        for earlier in created:
            await earlier.unload()
        provider = Qwen3AbstractionProvider(
            model_id=committed_qwen3_model_id() if model_id is None else model_id
        )
        created.append(provider)
        return provider

    with real_model_lock():
        try:
            yield make
        finally:
            for provider in created:
                release_provider(provider)
    for provider in created:
        _assert_released(provider)


@requires_real_models
@requires_qwen3
class TestQwen3AbstractionProvider:
    """Tests AD-026 through AD-033."""

    def test_ad_026_init_defers_loading(self):
        """AD-026: Provider constructor succeeds without loading model."""
        provider = Qwen3AbstractionProvider(model_id=committed_qwen3_model_id())
        assert provider._model is None
        assert provider._tokenizer is None

    def test_ad_026_bad_model_succeeds_at_init(self):
        """AD-026: Bad model ID succeeds at construction (lazy loading)."""
        provider = Qwen3AbstractionProvider(model_id="nonexistent-model-xyz")
        # Construction succeeds; failure deferred to first use
        assert provider._model is None

    async def test_ad_026_bad_model_fails_on_first_call(self):
        """AD-026: Bad model ID raises RuntimeError on first generate_abstract()."""
        provider = Qwen3AbstractionProvider(model_id="nonexistent-model-xyz")
        with pytest.raises(RuntimeError, match="nonexistent-model-xyz"):
            await provider.generate_abstract("Test text.", 200, None)

    async def test_ad_027_non_empty_output(self, qwen3_provider):
        """AD-027: Generated abstract is a non-empty string."""
        result = await qwen3_provider.generate_abstract(SAMPLE_TEXT, 200, None)
        assert isinstance(result, str)
        assert len(result.strip()) > 0
        assert result == result.strip()  # No leading/trailing whitespace

    async def test_ad_028_max_tokens_bound(self, qwen3_provider):
        """AD-028: Output respects max_tokens upper bound."""
        # Use a longer input to encourage full-length generation
        long_text = SAMPLE_TEXT * 5
        result = await qwen3_provider.generate_abstract(long_text, 100, None)

        # Tokenize the result using the model's tokenizer
        tokens = qwen3_provider._tokenizer.encode(result)
        assert len(tokens) <= 100, f"Abstract has {len(tokens)} tokens, expected at most 100"

    async def test_ad_029_deterministic(self, qwen3_provider):
        """AD-029: Same input produces identical output."""
        r1 = await qwen3_provider.generate_abstract(SAMPLE_TEXT, 200, None)
        r2 = await qwen3_provider.generate_abstract(SAMPLE_TEXT, 200, None)
        assert r1 == r2

    async def test_ad_030_short_input(self, qwen3_provider):
        """AD-030: Short input produces a valid abstract."""
        result = await qwen3_provider.generate_abstract(
            "Brief note about record linkage.", 200, None
        )
        assert isinstance(result, str)
        assert len(result.strip()) > 0

    async def test_ad_031_long_input(self, qwen3_provider):
        """AD-031: Long input does not crash."""
        # 50,000+ characters of repeated technical prose
        very_long_text = SAMPLE_TEXT * 200
        assert len(very_long_text) > 50_000

        result = await qwen3_provider.generate_abstract(very_long_text, 200, None)
        assert isinstance(result, str)
        assert len(result.strip()) > 0

    async def test_ad_032_semantic_quality(self, qwen3_provider, embedding_provider):
        """AD-032: Abstract is semantically related to input."""
        abstract = await qwen3_provider.generate_abstract(SAMPLE_TEXT, 200, None)

        # Embed both input and abstract
        vecs = await embedding_provider.embed([SAMPLE_TEXT, abstract])
        similarity = _cosine_sim(vecs[0], vecs[1])

        assert similarity > 0.5, f"cosine_similarity={similarity:.4f}, expected > 0.5"

        # Keyword overlap: at least one key concept appears in the abstract
        key_terms = ["health", "record", "clinical", "CRDT", "synchron"]
        abstract_lower = abstract.lower()
        matches = [t for t in key_terms if t.lower() in abstract_lower]
        assert len(matches) >= 1, f"No key terms found in abstract. Terms checked: {key_terms}"

    async def test_ad_033_error_propagation(self, qwen3_provider, monkeypatch):
        """AD-033: LLM runtime error propagates as exception."""

        def failing_generate(*args, **kwargs):
            raise RuntimeError("Simulated MLX inference failure")

        monkeypatch.setattr(qwen3_provider, "_generate_fn", failing_generate)

        with pytest.raises(RuntimeError, match="Simulated MLX inference failure"):
            await qwen3_provider.generate_abstract(SAMPLE_TEXT, 200, None)


@requires_real_models
@requires_qwen3
class TestQwen3LazyLoading:
    """Tests AD-095 through AD-097: lazy model loading behavior."""

    async def test_ad_095_defers_load_to_first_call(self, qwen3_provider_factory):
        """AD-095: Constructor does not load model; first call triggers load."""
        provider = await qwen3_provider_factory()
        assert provider._model is None

        result = await provider.generate_abstract(SAMPLE_TEXT, 200, None)
        assert provider._model is not None
        assert isinstance(result, str)
        assert len(result.strip()) > 0

    async def test_ad_096_second_call_reuses_model(self, qwen3_provider_factory):
        """AD-096: Second generate_abstract() reuses the loaded model."""
        provider = await qwen3_provider_factory()

        await provider.generate_abstract(SAMPLE_TEXT, 200, None)
        model_after_first = provider._model
        assert model_after_first is not None

        await provider.generate_abstract("Different input text.", 200, None)
        assert provider._model is model_after_first  # Same object identity

    async def test_ad_097_load_failure_raises_and_stays_unloaded(self, qwen3_provider_factory):
        """AD-097: Model load failure raises RuntimeError; provider
        remains in unloaded state for potential retry."""
        provider = await qwen3_provider_factory("nonexistent-model-xyz")

        with pytest.raises(RuntimeError, match="nonexistent-model-xyz"):
            await provider.generate_abstract(SAMPLE_TEXT, 200, None)

        assert provider._model is None  # No partial state


# ── Markdown Adapter: Source Provenance ────────────────────────────

import os  # noqa: E402 -- grouped with the markdown-adapter test section below
from datetime import (  # noqa: E402 -- grouped with the markdown-adapter test section below
    datetime,
    timezone,
)
from pathlib import Path  # noqa: E402 -- grouped with the markdown-adapter test section below

from sage.source_adapters.markdown_adapter import (  # noqa: E402 -- grouped with the markdown-adapter test section below
    MarkdownAdapter,
)


class TestMarkdownAdapterProvenance:
    """AD-034: Markdown adapter extracts source_modified_at from file mtime."""

    async def test_ad_034_source_modified_at_in_metadata(self, tmp_path):
        """AD-034: Markdown adapter populates source_modified_at from st_mtime."""
        adapter = MarkdownAdapter()

        test_file = tmp_path / "test_provenance.md"
        test_file.write_text("# Test\n\nContent for provenance test.")

        # Set a known mtime
        known_mtime = datetime(2023, 4, 20, 10, 0, 0, tzinfo=timezone.utc)
        os.utime(test_file, (test_file.stat().st_atime, known_mtime.timestamp()))

        result = await adapter.project(test_file)

        assert "source_modified_at" in result.metadata
        parsed = datetime.fromisoformat(result.metadata["source_modified_at"])
        assert parsed.tzinfo is not None  # timezone-aware
        assert abs((parsed - known_mtime).total_seconds()) < 1.0


class TestMarkdownAdapterCodeBlockSuppression:
    """ATX heading-shaped lines inside code blocks must not extract.

    The regex parser at adapter_version 0.3.0 had no awareness of fenced or
    indented code-block context, so heading-format examples inside them were
    indexed as real headings and corrupted the heading hierarchy.
    """

    async def test_tc1_fenced_heading_not_extracted(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc1.md"
        test_file.write_text("# Real\n\n```\n## Phantom\n```\n")

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert "Phantom" not in texts
        assert texts == ["Real"]
        assert result.headings[0].level == 1

    async def test_tc2_indented_code_block_heading_not_extracted(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc2.md"
        test_file.write_text("# Real\n\nbody\n\n    ## Phantom\n")

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert "Phantom" not in texts
        assert texts == ["Real"]

    async def test_tc3_language_tagged_fence_heading_not_extracted(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc3.md"
        test_file.write_text("# Real\n\n```python\n## Phantom\n```\n")

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert "Phantom" not in texts
        assert texts == ["Real"]

    async def test_tc4_all_six_levels_in_fence_suppressed(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc4.md"
        test_file.write_text(
            "# Real\n\n```\n# h1\n## h2\n### h3\n#### h4\n##### h5\n###### h6\n```\n"
        )

        result = await adapter.project(test_file)

        assert len(result.headings) == 1
        assert result.headings[0].text == "Real"
        assert result.headings[0].level == 1

    async def test_tc5_hierarchy_intact_across_phantom(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc5.md"
        test_file.write_text("# A\n\n## B\n\n```\n## Phantom\n```\n\n## C\n")

        result = await adapter.project(test_file)

        by_text = {h.text: h for h in result.headings}
        assert "Phantom" not in by_text
        assert by_text["B"].level == 2
        assert by_text["C"].level == 2
        assert by_text["B"].path == "A > B"
        assert by_text["C"].path == "A > C"

    async def test_tc6_canonical_reproducer_open_questions_resolved(self, tmp_path):
        """Mirrors the example_vault design-note bug where `## Open Questions
        Resolved` inside a fenced block phantom-headed §3.2..§3.5."""
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc6.md"
        test_file.write_text(
            "# Doc\n\n## 3\n\n### 3.1\n\n```\n## Open Questions Resolved\n```\n\n### 3.2\n"
        )

        result = await adapter.project(test_file)

        by_text = {h.text: h for h in result.headings}
        assert "Open Questions Resolved" not in by_text
        assert by_text["3.2"].path == "Doc > 3 > 3.2"
        assert by_text["3.2"].level == 3


class TestMarkdownAdapterFrontmatterStripping:
    """YAML frontmatter must not extract as headings.

    Without the front_matter_plugin, the CommonMark parser binds the YAML
    body as a paragraph anchored to its closing ``---`` delimiter and
    interprets that delimiter as a setext H2 underline, admitting the raw
    YAML payload as the first H2 in the heading index.
    """

    async def test_tc1_multiline_description_no_yaml_in_headings(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc1.md"
        test_file.write_text(
            "---\n"
            "name: skill-name\n"
            "description: A multi-line description that\n"
            "  wraps onto a second line via YAML block scalar.\n"
            "---\n"
            "\n"
            "# Real H1\n"
            "\n"
            "## Real H2\n"
        )

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert texts == ["Real H1", "Real H2"]
        assert result.headings[0].level == 1
        assert result.headings[1].level == 2
        assert not any("name:" in t or "description:" in t for t in texts)

    async def test_tc2_single_line_description_no_yaml_in_headings(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc2.md"
        test_file.write_text(
            "---\nname: skill\ndescription: One-line description.\n---\n\n# Real H1\n"
        )

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert texts == ["Real H1"]
        assert not any("name:" in t or "description:" in t for t in texts)

    async def test_tc3_quoted_multiline_description_with_version(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc3.md"
        test_file.write_text(
            "---\n"
            'description: "A quoted multi-line\n'
            '  description that wraps."\n'
            "version: 1.2.3\n"
            "---\n"
            "\n"
            "# Real H1\n"
            "\n"
            "## Real H2\n"
        )

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert texts == ["Real H1", "Real H2"]
        assert not any("description:" in t or "version:" in t for t in texts)

    async def test_tc4_no_frontmatter_unchanged(self, tmp_path):
        """Regression guard: documents without frontmatter parse identically
        to adapter_version 0.4.0 behavior."""
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc4.md"
        test_file.write_text("# H1\n\n## H2\n\nbody\n\n### H3\n")

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert texts == ["H1", "H2", "H3"]
        assert result.headings[0].level == 1
        assert result.headings[1].level == 2
        assert result.headings[2].level == 3

    async def test_tc5a_doc_start_dash_with_blank_is_thematic_break(self, tmp_path):
        """A document opening with ``---`` followed by a blank line and
        no later ``---``/``...`` closer is not a frontmatter opener.
        The plugin's end-marker search fails, the rule aborts, and the
        standard CommonMark thematic-break parser handles the opening
        ``---`` instead."""
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc5a.md"
        test_file.write_text("---\n\n# H1\n\nbody paragraph with no body-internal `---`.\n")

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert texts == ["H1"]
        assert result.headings[0].level == 1

    async def test_tc5b_body_internal_thematic_break_unaffected(self, tmp_path):
        """A body-internal ``---`` thematic break in a document that does
        not start with ``---`` is never a frontmatter candidate (the
        plugin only fires at line 0). Both H1 and the H2 following the
        thematic break are extracted."""
        adapter = MarkdownAdapter()
        test_file = tmp_path / "tc5b.md"
        test_file.write_text("# H1\n\nbody paragraph.\n\n---\n\n## H2\n")

        result = await adapter.project(test_file)

        texts = [h.text for h in result.headings]
        assert texts == ["H1", "H2"]


class TestMarkdownAdapterADRTier3Extraction:
    """Markdown adapter extracts ``adr_id`` from ``cas-adr-NNN_*`` filenames.

    The cas vault names ADR sources with the lowercase ``cas-adr-NNN_<title>``
    convention. Per CAS-ADR-021, the adapter surfaces this filename-derived
    tier3 fact via ``ProjectionResult.metadata["adapter_tier3_metadata"]``.
    The ingestion service merges it below caller-supplied tier3 (caller wins).
    """

    async def test_adr_filename_emits_adr_id_in_adapter_tier3(self, tmp_path):
        adapter = MarkdownAdapter()
        test_file = tmp_path / "cas-adr-038_Some_Title.md"
        test_file.write_text("# ADR-038: Some Title\n\nBody.\n")

        result = await adapter.project(test_file)

        assert result.metadata.get("adapter_tier3_metadata") == {"adr_id": "038"}

    async def test_non_adr_markdown_emits_no_adapter_tier3(self, tmp_path):
        """Anti-coincidence: only ``cas-adr-NNN_*`` filenames emit tier3."""
        adapter = MarkdownAdapter()

        date_prefixed = tmp_path / "2026-05-27_CAS_Some_Note.md"
        date_prefixed.write_text("# Note\n\nBody.\n")
        result_a = await adapter.project(date_prefixed)
        assert "adapter_tier3_metadata" not in result_a.metadata

        bare_name = tmp_path / "not-an-adr_file.md"
        bare_name.write_text("# Title\n\nBody.\n")
        result_b = await adapter.project(bare_name)
        assert "adapter_tier3_metadata" not in result_b.metadata

    @pytest.mark.parametrize(
        "filename",
        [
            "cas-adr-9_Title.md",  # too few digits
            "cas-adr-9999_Title.md",  # too many digits
            "CAS-ADR-099_Title.md",  # uppercase
            "cas-adr-099.md",  # no trailing title segment
        ],
    )
    async def test_malformed_adr_filenames_emit_no_adapter_tier3(self, tmp_path, filename):
        """The extraction regex rejects shapes that violate the convention.

        The convention is strict: lowercase ``cas-adr-`` prefix, exactly
        three digits, an underscore separator, then a title segment.
        Anything else is not an ADR for extraction purposes — the adapter
        emits no tier3 hint and the ingestion path falls back to the
        caller (or to no tier3 at all).
        """
        adapter = MarkdownAdapter()
        test_file = tmp_path / filename
        test_file.write_text("# Title\n\nBody.\n")

        result = await adapter.project(test_file)

        assert "adapter_tier3_metadata" not in result.metadata


# ── Docx Adapter ────────────────────────────────────────────────────

import hashlib  # noqa: E402 -- grouped with the docx-adapter test section below

try:
    import docx
    from docx.oxml.ns import qn
    from lxml import etree

    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

requires_docx = pytest.mark.skipif(not _HAS_DOCX, reason="python-docx not available")


def _make_docx(tmp_path: Path, filename: str = "test.docx") -> Path:
    """Create a minimal empty.docx and return its path."""
    doc = docx.Document()
    path = tmp_path / filename
    doc.save(str(path))
    return path


def _add_heading_with_style(doc, text: str, style_name: str) -> None:
    """Add a paragraph with the given built-in heading style."""
    doc.add_paragraph(text, style=style_name)


def _add_table(doc, rows: list[list[str]]) -> None:
    """Add a table with the given cell data."""
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            table.cell(i, j).text = cell_text


def _inject_numbering(doc, abstract_num_xml: str, num_xml: str) -> None:
    """Inject numbering definitions into the document's numbering part.

    Creates the numbering part if it doesn't exist, then appends the
    provided abstractNum and num elements.
    """
    # Ensure numbering part exists by adding and removing a dummy list
    doc.add_paragraph("dummy", style="List Bullet")
    # Remove the dummy paragraph
    body = doc.element.body
    body.remove(body[-1])

    numbering_part = doc.part.numbering_part
    numbering_elm = numbering_part.numbering_definitions._numbering

    # Parse and append abstractNum
    abstract_elem = etree.fromstring(abstract_num_xml)
    numbering_elm.append(abstract_elem)

    # Parse and append num
    num_elem = etree.fromstring(num_xml)
    numbering_elm.append(num_elem)


def _set_paragraph_numbering(paragraph, num_id: int, ilvl: int) -> None:
    """Set w:numPr on a paragraph to reference a numbering definition."""
    pPr = paragraph._element.get_or_add_pPr()
    numPr = etree.SubElement(pPr, qn("w:numPr"))
    ilvl_elem = etree.SubElement(numPr, qn("w:ilvl"))
    ilvl_elem.set(qn("w:val"), str(ilvl))
    numId_elem = etree.SubElement(numPr, qn("w:numId"))
    numId_elem.set(qn("w:val"), str(num_id))


def _inject_cross_ref_field(paragraph, instruction: str, display_text: str) -> None:
    """Inject a complex field (fldChar begin/separate/end) into a paragraph.

    Simulates a cross-reference field like REF _Ref12345 \\r \\h with a
    cached display value.
    """
    p_elem = paragraph._element

    # Begin
    r_begin = etree.SubElement(p_elem, qn("w:r"))
    fc_begin = etree.SubElement(r_begin, qn("w:fldChar"))
    fc_begin.set(qn("w:fldCharType"), "begin")

    # Instruction
    r_instr = etree.SubElement(p_elem, qn("w:r"))
    instr_text = etree.SubElement(r_instr, qn("w:instrText"))
    instr_text.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    instr_text.text = instruction

    # Separate
    r_sep = etree.SubElement(p_elem, qn("w:r"))
    fc_sep = etree.SubElement(r_sep, qn("w:fldChar"))
    fc_sep.set(qn("w:fldCharType"), "separate")

    # Display text (cached result)
    r_display = etree.SubElement(p_elem, qn("w:r"))
    t_display = etree.SubElement(r_display, qn("w:t"))
    t_display.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    t_display.text = display_text

    # End
    r_end = etree.SubElement(p_elem, qn("w:r"))
    fc_end = etree.SubElement(r_end, qn("w:fldChar"))
    fc_end.set(qn("w:fldCharType"), "end")


# Standard decimal numbering XML for heading levels 0-2
DECIMAL_ABSTRACT_NUM_XML = """
<w:abstractNum w:abstractNumId="100"
    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:lvl w:ilvl="0">
    <w:start w:val="1"/>
    <w:numFmt w:val="decimal"/>
    <w:lvlText w:val="%1"/>
    <w:lvlJc w:val="left"/>
  </w:lvl>
  <w:lvl w:ilvl="1">
    <w:start w:val="1"/>
    <w:numFmt w:val="decimal"/>
    <w:lvlText w:val="%1.%2"/>
    <w:lvlJc w:val="left"/>
  </w:lvl>
  <w:lvl w:ilvl="2">
    <w:start w:val="1"/>
    <w:numFmt w:val="decimal"/>
    <w:lvlText w:val="%1.%2.%3"/>
    <w:lvlJc w:val="left"/>
  </w:lvl>
</w:abstractNum>
"""

DECIMAL_NUM_XML = """
<w:num w:numId="100"
    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:abstractNumId w:val="100"/>
</w:num>
"""

# Mixed numbering: Roman, decimal, lowercase alpha
MIXED_ABSTRACT_NUM_XML = """
<w:abstractNum w:abstractNumId="200"
    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:lvl w:ilvl="0">
    <w:start w:val="1"/>
    <w:numFmt w:val="upperRoman"/>
    <w:lvlText w:val="%1"/>
    <w:lvlJc w:val="left"/>
  </w:lvl>
  <w:lvl w:ilvl="1">
    <w:start w:val="1"/>
    <w:numFmt w:val="decimal"/>
    <w:lvlText w:val="%1.%2"/>
    <w:lvlJc w:val="left"/>
  </w:lvl>
  <w:lvl w:ilvl="2">
    <w:start w:val="1"/>
    <w:numFmt w:val="lowerLetter"/>
    <w:lvlText w:val="%1.%2.%3"/>
    <w:lvlJc w:val="left"/>
  </w:lvl>
</w:abstractNum>
"""

MIXED_NUM_XML = """
<w:num w:numId="200"
    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:abstractNumId w:val="200"/>
</w:num>
"""


@requires_docx
class TestDocxAdapter:
    """AD-035 through AD-052: Docx source adapter tests."""

    async def test_ad_035_basic_projection(self, tmp_path):
        """AD-035: Basic projection returns valid ProjectionResult."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Introduction", style="Heading 1")
        doc.add_paragraph("This is the introduction text.")
        path = tmp_path / "basic.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert isinstance(result.text, str)
        assert len(result.text) > 0
        assert len(result.headings) == 1
        assert result.headings[0].text == "Introduction"
        assert isinstance(result.content_hash, str)
        assert len(result.content_hash) == 64  # SHA-256 hex
        assert result.title == "basic"  # filename stem (no Title-styled paragraph)
        assert result.adapter_version == DocxAdapter.VERSION

    async def test_ad_036_heading_style_map_config(self, tmp_path):
        """AD-036: Heading extraction uses heading_style_map config."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("My Title", style="Title")
        doc.add_paragraph("Body under title.")
        path = tmp_path / "custom_style.docx"
        doc.save(str(path))

        config = {"heading_style_map": {"Title": 1}}
        result = await adapter.project(path, config=config)

        assert len(result.headings) == 1
        assert result.headings[0].level == 1
        assert result.headings[0].text == "My Title"

    async def test_ad_037_default_heading_styles(self, tmp_path):
        """AD-037: Default heading styles work without explicit config."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Level One", style="Heading 1")
        doc.add_paragraph("Level Two", style="Heading 2")
        doc.add_paragraph("Level Three", style="Heading 3")
        path = tmp_path / "defaults.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert len(result.headings) == 3
        assert result.headings[0].level == 1
        assert result.headings[1].level == 2
        assert result.headings[2].level == 3

    async def test_ad_038_heading_hierarchy_paths(self, tmp_path):
        """AD-038: Heading hierarchy paths use ' > ' separator."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Chapter", style="Heading 1")
        doc.add_paragraph("Section", style="Heading 2")
        doc.add_paragraph("Subsection", style="Heading 3")
        path = tmp_path / "hierarchy.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert result.headings[0].path == "Chapter"
        assert result.headings[1].path == "Chapter > Section"
        assert result.headings[2].path == "Chapter > Section > Subsection"

    async def test_ad_039_title_extraction_priority_chain(self, tmp_path):
        """AD-039: Title priority: Title style > filename > key terms."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()

        # Priority 1: Title paragraph style takes precedence over H1
        doc1 = docx.Document()
        doc1.add_paragraph("Formal Document Title", style="Title")
        doc1.add_paragraph("System Architecture", style="Heading 1")
        path1 = tmp_path / "AuthoritativeAccumulator.docx"
        doc1.save(str(path1))
        result1 = await adapter.project(path1)
        assert result1.title == "Formal Document Title"

        # Priority 2: Filename stem when no Title style present
        doc2 = docx.Document()
        doc2.add_paragraph("Introduction", style="Heading 1")
        doc2.add_paragraph("Some body text about clinical normalization.")
        path2 = tmp_path / "AuthoritativeAccumulator.docx"
        doc2.save(str(path2))
        result2 = await adapter.project(path2)
        assert result2.title == "AuthoritativeAccumulator"

        # Priority 3: Key terms from first body paragraph when no Title
        # style and filename is degenerate
        doc3 = docx.Document()
        doc3.add_paragraph("Abstract", style="Heading 1")
        doc3.add_paragraph(
            "The authoritative accumulator validates clinical "
            "normalization through respiratory signal processing."
        )
        path3 = tmp_path / ".docx"
        doc3.save(str(path3))
        result3 = await adapter.project(path3)
        # Should contain content words, not stop words
        assert "authoritative" in result3.title.lower()
        assert "the" not in result3.title.lower().split()

    async def test_ad_040_content_hash_is_raw_bytes(self, tmp_path):
        """AD-040: content_hash is SHA-256 of raw.docx bytes."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Hash test content.")
        path = tmp_path / "hash_test.docx"
        doc.save(str(path))

        expected_hash = hashlib.sha256(path.read_bytes()).hexdigest()
        result = await adapter.project(path)
        assert result.content_hash == expected_hash

    async def test_ad_041_source_modified_at(self, tmp_path):
        """AD-041: source_modified_at extracted from file mtime."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Provenance test.")
        path = tmp_path / "provenance.docx"
        doc.save(str(path))

        known_mtime = datetime(2023, 6, 15, 12, 0, 0, tzinfo=timezone.utc)
        os.utime(path, (path.stat().st_atime, known_mtime.timestamp()))

        result = await adapter.project(path)

        assert "source_modified_at" in result.metadata
        parsed = datetime.fromisoformat(result.metadata["source_modified_at"])
        assert parsed.tzinfo is not None
        assert abs((parsed - known_mtime).total_seconds()) < 1.0

    async def test_ad_042_table_extraction(self, tmp_path):
        """AD-042: Table content extracted as pipe-delimited text rows."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Data Section", style="Heading 1")
        _add_table(doc, [["Name", "Value"], ["alpha", "1"], ["beta", "2"]])
        path = tmp_path / "table.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert "| Name | Value |" in result.text
        assert "| alpha | 1 |" in result.text
        assert "| beta | 2 |" in result.text

    async def test_ad_043_mixed_content_ordering(self, tmp_path):
        """AD-043: Mixed headings, paragraphs, tables in correct order."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Overview", style="Heading 1")
        doc.add_paragraph("Intro paragraph.")
        _add_table(doc, [["A", "B"], ["1", "2"]])
        doc.add_paragraph("Details", style="Heading 2")
        doc.add_paragraph("Detail paragraph.")
        path = tmp_path / "mixed.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert len(result.headings) == 2
        assert result.headings[0].text == "Overview"
        assert result.headings[1].text == "Details"
        # Table content should be under "Overview" heading
        assert "| A | B |" in result.headings[0].content
        assert "Intro paragraph." in result.headings[0].content

    async def test_ad_044_empty_document(self, tmp_path):
        """AD-044: Empty document produces valid result with filename title."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        path = tmp_path / "empty_doc.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert result.title == "empty_doc"
        assert result.headings == []
        assert isinstance(result.text, str)
        assert isinstance(result.content_hash, str)

    async def test_ad_045_custom_style_map_override(self, tmp_path):
        """AD-045: Custom heading_style_map overrides defaults."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Custom Top", style="Title")
        doc.add_paragraph("Standard H1", style="Heading 1")
        path = tmp_path / "override.docx"
        doc.save(str(path))

        # "Title" mapped to level 1; defaults still apply for "Heading 1"
        config = {"heading_style_map": {"Title": 1}}
        result = await adapter.project(path, config=config)

        titles = [h for h in result.headings if h.text == "Custom Top"]
        assert len(titles) == 1
        assert titles[0].level == 1

        h1s = [h for h in result.headings if h.text == "Standard H1"]
        assert len(h1s) == 1
        assert h1s[0].level == 1

    async def test_ad_046_body_under_heading(self, tmp_path):
        """AD-046: Non-heading paragraphs appear as content under nearest heading."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Section A", style="Heading 1")
        doc.add_paragraph("First body paragraph.")
        doc.add_paragraph("Second body paragraph.")
        doc.add_paragraph("Section B", style="Heading 1")
        doc.add_paragraph("Third body paragraph.")
        path = tmp_path / "body.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert "First body paragraph." in result.headings[0].content
        assert "Second body paragraph." in result.headings[0].content
        assert "Third body paragraph." in result.headings[1].content
        assert "First body paragraph." not in result.headings[1].content

    async def test_ad_047_adapter_version(self, tmp_path):
        """AD-047: adapter_version matches DocxAdapter.VERSION."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Version check.")
        path = tmp_path / "version.docx"
        doc.save(str(path))

        result = await adapter.project(path)
        assert result.adapter_version == DocxAdapter.VERSION

    async def test_ad_048_decimal_heading_numbering(self, tmp_path):
        """AD-048: Decimal heading numbering prepended to heading text."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        _inject_numbering(doc, DECIMAL_ABSTRACT_NUM_XML, DECIMAL_NUM_XML)

        # Two H1 headings, second with an H2 child
        p1 = doc.add_paragraph("Introduction", style="Heading 1")
        _set_paragraph_numbering(p1, num_id=100, ilvl=0)

        p2 = doc.add_paragraph("Background", style="Heading 1")
        _set_paragraph_numbering(p2, num_id=100, ilvl=0)

        p3 = doc.add_paragraph("Definitions", style="Heading 2")
        _set_paragraph_numbering(p3, num_id=100, ilvl=1)

        path = tmp_path / "numbered.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert result.headings[0].text == "1 Introduction"
        assert result.headings[1].text == "2 Background"
        assert result.headings[2].text == "2.1 Definitions"

    async def test_ad_049_numbering_counter_reset(self, tmp_path):
        """AD-049: Child counters reset when parent level increments."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        _inject_numbering(doc, DECIMAL_ABSTRACT_NUM_XML, DECIMAL_NUM_XML)

        # H1 -> H2, H2 -> H1 -> H2 (second H2 should restart at.1)
        p1 = doc.add_paragraph("Part A", style="Heading 1")
        _set_paragraph_numbering(p1, num_id=100, ilvl=0)
        p2 = doc.add_paragraph("Sub One", style="Heading 2")
        _set_paragraph_numbering(p2, num_id=100, ilvl=1)
        p3 = doc.add_paragraph("Sub Two", style="Heading 2")
        _set_paragraph_numbering(p3, num_id=100, ilvl=1)
        p4 = doc.add_paragraph("Part B", style="Heading 1")
        _set_paragraph_numbering(p4, num_id=100, ilvl=0)
        p5 = doc.add_paragraph("Sub One Again", style="Heading 2")
        _set_paragraph_numbering(p5, num_id=100, ilvl=1)

        path = tmp_path / "reset.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert result.headings[0].text == "1 Part A"
        assert result.headings[1].text == "1.1 Sub One"
        assert result.headings[2].text == "1.2 Sub Two"
        assert result.headings[3].text == "2 Part B"
        assert result.headings[4].text == "2.1 Sub One Again"

    async def test_ad_050_custom_numbering_formats(self, tmp_path):
        """AD-050: upperRoman and lowerLetter numbering formats."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        _inject_numbering(doc, MIXED_ABSTRACT_NUM_XML, MIXED_NUM_XML)

        p1 = doc.add_paragraph("First Chapter", style="Heading 1")
        _set_paragraph_numbering(p1, num_id=200, ilvl=0)
        p2 = doc.add_paragraph("First Section", style="Heading 2")
        _set_paragraph_numbering(p2, num_id=200, ilvl=1)
        p3 = doc.add_paragraph("First Item", style="Heading 3")
        _set_paragraph_numbering(p3, num_id=200, ilvl=2)

        path = tmp_path / "roman.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert result.headings[0].text == "I First Chapter"
        assert result.headings[1].text == "I.1 First Section"
        assert result.headings[2].text == "I.1.a First Item"

    async def test_ad_051_cross_reference_field(self, tmp_path):
        """AD-051: Cross-ref field cached results in text, instructions excluded."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("References", style="Heading 1")

        # Add paragraph with text + cross-reference field
        p = doc.add_paragraph()
        p.add_run("See Section ")
        _inject_cross_ref_field(p, " REF _Ref12345 \\r \\h ", "1.1")
        p.add_run(" for details.")

        path = tmp_path / "crossref.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        # Cached display value should appear
        assert "See Section 1.1 for details." in result.text
        # Field instruction should NOT appear
        assert "REF _Ref12345" not in result.text
        assert "instrText" not in result.text

    async def test_ad_052_mixed_numbering_formats(self, tmp_path):
        """AD-052: Multi-level numbering with mixed formats (I.2.a)."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        _inject_numbering(doc, MIXED_ABSTRACT_NUM_XML, MIXED_NUM_XML)

        # Two top-level chapters, each with subsections
        p1 = doc.add_paragraph("Alpha", style="Heading 1")
        _set_paragraph_numbering(p1, num_id=200, ilvl=0)
        p2 = doc.add_paragraph("Sub Alpha", style="Heading 2")
        _set_paragraph_numbering(p2, num_id=200, ilvl=1)
        p3 = doc.add_paragraph("Detail A", style="Heading 3")
        _set_paragraph_numbering(p3, num_id=200, ilvl=2)
        p4 = doc.add_paragraph("Detail B", style="Heading 3")
        _set_paragraph_numbering(p4, num_id=200, ilvl=2)
        p5 = doc.add_paragraph("Beta", style="Heading 1")
        _set_paragraph_numbering(p5, num_id=200, ilvl=0)
        p6 = doc.add_paragraph("Sub Beta", style="Heading 2")
        _set_paragraph_numbering(p6, num_id=200, ilvl=1)

        path = tmp_path / "mixed_fmt.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert result.headings[0].text == "I Alpha"
        assert result.headings[1].text == "I.1 Sub Alpha"
        assert result.headings[2].text == "I.1.a Detail A"
        assert result.headings[3].text == "I.1.b Detail B"
        assert result.headings[4].text == "II Beta"
        assert result.headings[5].text == "II.1 Sub Beta"


@requires_docx
class TestDocxAdapterADRTier3Extraction:
    """Docx adapter extracts ``adr_id`` from ``cas-adr-NNN_*`` filenames.

    Mirrors the markdown adapter behaviour. The cas vault's ADR doc_type
    declares ``source_types: [docx, markdown]``; both adapters share the
    helper at ``sage.source_adapters.base.extract_adr_id_from_filename``.
    """

    async def test_adr_docx_filename_emits_adr_id_in_adapter_tier3(self, tmp_path):
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("ADR-038: Some Title", style="Heading 1")
        doc.add_paragraph("Body.")
        path = tmp_path / "cas-adr-038_Some_Title.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert result.metadata.get("adapter_tier3_metadata") == {"adr_id": "038"}

    async def test_non_adr_docx_emits_no_adapter_tier3(self, tmp_path):
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Some Note", style="Heading 1")
        path = tmp_path / "2026-05-27_CAS_Some_Note.docx"
        doc.save(str(path))

        result = await adapter.project(path)

        assert "adapter_tier3_metadata" not in result.metadata

    @pytest.mark.parametrize(
        "filename",
        [
            "cas-adr-9_Title.docx",  # too few digits
            "cas-adr-9999_Title.docx",  # too many digits
            "CAS-ADR-099_Title.docx",  # uppercase
            "cas-adr-099.docx",  # no trailing title segment
        ],
    )
    async def test_malformed_adr_docx_filenames_emit_no_adapter_tier3(self, tmp_path, filename):
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        doc = docx.Document()
        doc.add_paragraph("Title", style="Heading 1")
        path = tmp_path / filename
        doc.save(str(path))

        result = await adapter.project(path)

        assert "adapter_tier3_metadata" not in result.metadata


# ── Docx Adapter:.dotx template support ──────────────────────────


def _add_custom_style(
    doc,
    style_id: str,
    style_name: str,
    style_type: str = "paragraph",
    based_on: str | None = None,
    num_id: int | None = None,
) -> None:
    """Inject a <w:style> element into styles.xml with customStyle="1".

    python-docx does not expose `customStyle` when adding styles via its
    API, so we go directly to XML. `num_id` links a paragraph style to a
    numbering definition created via _inject_numbering().
    """
    style_type_attr = {
        "paragraph": "paragraph",
        "character": "character",
        "table": "table",
        "numbering": "numbering",
    }[style_type]
    style_elem = etree.SubElement(doc.styles.element, qn("w:style"))
    style_elem.set(qn("w:type"), style_type_attr)
    style_elem.set(qn("w:styleId"), style_id)
    style_elem.set(qn("w:customStyle"), "1")
    name_elem = etree.SubElement(style_elem, qn("w:name"))
    name_elem.set(qn("w:val"), style_name)
    if based_on is not None:
        based_on_elem = etree.SubElement(style_elem, qn("w:basedOn"))
        based_on_elem.set(qn("w:val"), based_on)
    if num_id is not None and style_type == "paragraph":
        pPr = etree.SubElement(style_elem, qn("w:pPr"))
        numPr = etree.SubElement(pPr, qn("w:numPr"))
        ilvl = etree.SubElement(numPr, qn("w:ilvl"))
        ilvl.set(qn("w:val"), "0")
        numId = etree.SubElement(numPr, qn("w:numId"))
        numId.set(qn("w:val"), str(num_id))


def _attach_numbering_to_builtin_style(doc, style_id: str, num_id: int) -> None:
    """Wire an existing style (custom or built-in) to a numbering definition.

    Ensures the style's w:pPr contains a numPr referencing num_id. Used to
    simulate a template that has customized the built-in Heading 1 style
    with auto-numbering.
    """
    style = doc.styles[style_id]
    pPr = style.element.find(qn("w:pPr"))
    if pPr is None:
        pPr = etree.SubElement(style.element, qn("w:pPr"))
    # Remove any existing numPr
    existing = pPr.find(qn("w:numPr"))
    if existing is not None:
        pPr.remove(existing)
    numPr = etree.SubElement(pPr, qn("w:numPr"))
    ilvl = etree.SubElement(numPr, qn("w:ilvl"))
    ilvl.set(qn("w:val"), "0")
    numId = etree.SubElement(numPr, qn("w:numId"))
    numId.set(qn("w:val"), str(num_id))


def _convert_docx_to_dotx(docx_path: Path, dotx_path: Path) -> None:
    """Rewrite a.docx's content-type entry to template flavor and save as.dotx.

    Creates a structurally valid.dotx (OPC content type set to template)
    that python-docx will reject at load unless the adapter's workaround
    applies, so this fixture exercises the real format path.
    """
    import zipfile as _zipfile

    with _zipfile.ZipFile(docx_path, "r") as z_in:
        with _zipfile.ZipFile(dotx_path, "w", _zipfile.ZIP_DEFLATED) as z_out:
            for item in z_in.namelist():
                data = z_in.read(item)
                if item == "[Content_Types].xml":
                    data = data.replace(
                        b"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
                        b"application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml",
                    )
                z_out.writestr(item, data)


def _build_template_fixture(tmp_path: Path, filename: str) -> Path:
    """Build a.dotx fixture with the style surface required by AD-071/073.

    The fixture contains:
    - AppendixHeading (custom, paragraph, no numbering)
    - DefinitionsHeader (custom, paragraph, active numbering)
    - AnnotationChar (custom, character, based_on Normal)
    - Built-in Heading 1 with active numbering (template-local wiring)
    - Built-in Normal unmodified
    """
    doc = docx.Document()
    doc.add_paragraph("Intro", style="Heading 1")
    doc.add_paragraph("Body paragraph.")

    # Numbering definitions for the two styles that need active numbering
    _inject_numbering(doc, DECIMAL_ABSTRACT_NUM_XML, DECIMAL_NUM_XML)
    second_num_xml = """
    <w:num w:numId="101"
        xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
      <w:abstractNumId w:val="100"/>
    </w:num>
    """
    numbering_part = doc.part.numbering_part
    numbering_elm = numbering_part.numbering_definitions._numbering
    numbering_elm.append(etree.fromstring(second_num_xml))

    _add_custom_style(doc, "AppendixHeading", "Appendix Heading", "paragraph")
    _add_custom_style(
        doc,
        "DefinitionsHeader",
        "Definitions Header",
        "paragraph",
        num_id=100,
    )
    _add_custom_style(
        doc,
        "AnnotationChar",
        "Annotation Char",
        "character",
        based_on="Normal",
    )
    _attach_numbering_to_builtin_style(doc, "Heading 1", num_id=101)

    docx_path = tmp_path / (filename.removesuffix(".dotx") + ".docx")
    dotx_path = tmp_path / filename
    doc.save(str(docx_path))
    _convert_docx_to_dotx(docx_path, dotx_path)
    docx_path.unlink()
    return dotx_path


@requires_docx
class TestDocxAdapterDotxSupport:
    """AD-068 through AD-073:.dotx template support in DocxAdapter."""

    async def test_ad_068_extensions_includes_dotx(self):
        """AD-068: DocxAdapter.EXTENSIONS contains both.docx and.dotx."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        assert ".docx" in DocxAdapter.EXTENSIONS
        assert ".dotx" in DocxAdapter.EXTENSIONS

    async def test_ad_069_dotx_loads_without_error(self, tmp_path):
        """AD-069: A.dotx file is parsed successfully despite template content type."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        dotx_path = _build_template_fixture(tmp_path, "fixture.dotx")

        result = await adapter.project(dotx_path)

        assert len(result.text) > 0
        assert isinstance(result.content_hash, str)
        assert len(result.content_hash) == 64
        assert len(result.headings) >= 1
        # content_hash is over raw.dotx bytes, not the shadow copy
        assert result.content_hash == hashlib.sha256(dotx_path.read_bytes()).hexdigest()

    async def test_ad_074_unreadable_docx_and_dotx_name_the_source_not_the_shadow(self, tmp_path):
        """A .docx or .dotx the library cannot open fails as a ValueError naming
        the file the caller supplied -- in the library's own text as well as in
        the prefix.

        Two properties. The library's exception (`PackageNotFoundError`, which
        is not a `ValueError`) is wrapped, so a corrupt document fails in the
        same shape the pdf and pptx adapters produce. And the .dotx branch
        opens a *shadow* copy this method writes under a temp directory, so the
        library names the shadow rather than the caller's file -- a scratch
        location that describes nothing the caller could act on.

        Anti-coincidental-pass: the second fixture is the load-bearing one, and
        it is built to reach the single python-docx failure that names the file
        it was handed on the shadow branch -- the content-type check, which
        fires only after the package opens cleanly. A malformed package cannot
        reach it: python-docx dies earlier with an error naming nothing, so a
        `sage_dotx_` assertion over that input is satisfied by any
        implementation, including one that interpolates the shadow verbatim.
        Hence a real package with its main-part content type rewritten to the
        macro-enabled-template type and a .dotx suffix, which takes the shadow
        branch (keyed on suffix) and then fails the content-type check against
        the shadow's path. Dropping the substitution turns this red; dropping
        the prefix turns the source assertion red.
        """
        import zipfile as _zipfile

        from sage.source_adapters.docx_adapter import _DOCX_CONTENT_TYPE, DocxAdapter

        adapter = DocxAdapter()

        not_a_zip = tmp_path / "corrupt.docx"
        not_a_zip.write_bytes(b"not a zip at all")
        with pytest.raises(ValueError) as excinfo:
            await adapter.project(not_a_zip)
        assert str(not_a_zip) in str(excinfo.value), excinfo.value

        # A well-formed package whose main part is typed as a macro-enabled
        # template: the shadow rewrite only swaps the plain dotx type, so this
        # reaches python-docx's content-type check with the shadow's path.
        built = tmp_path / "built.docx"
        docx.Document().save(str(built))
        macro_type = "application/vnd.ms-word.template.macroEnabledTemplate.main+xml"
        bad_template = tmp_path / "macro_template.dotx"
        with _zipfile.ZipFile(built) as z_in:
            with _zipfile.ZipFile(bad_template, "w", _zipfile.ZIP_DEFLATED) as z_out:
                for item in z_in.namelist():
                    data = z_in.read(item)
                    if item == "[Content_Types].xml":
                        data = data.replace(
                            _DOCX_CONTENT_TYPE.encode("utf-8"), macro_type.encode("utf-8")
                        )
                    z_out.writestr(item, data)

        with pytest.raises(ValueError) as excinfo:
            await adapter.project(bad_template)
        message = str(excinfo.value)
        # The precondition the whole test rests on: this input really did reach
        # the content-type failure, the only one that names a path here.
        assert "is not a Word file" in message, message
        assert str(bad_template) in message, message
        assert "sage_dotx_" not in message, message
        assert "shadow.docx" not in message, message

    async def test_ad_070_dotx_has_inventory_docx_does_not(self, tmp_path):
        """AD-070: template_style_inventory is.dotx-only; absent on.docx."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        dotx_path = _build_template_fixture(tmp_path, "t.dotx")

        # Equivalent.docx with the same body content and custom styles
        doc = docx.Document()
        doc.add_paragraph("Intro", style="Heading 1")
        doc.add_paragraph("Body paragraph.")
        _add_custom_style(doc, "AppendixHeading", "Appendix Heading", "paragraph")
        docx_path = tmp_path / "t.docx"
        doc.save(str(docx_path))

        dotx_result = await adapter.project(dotx_path)
        docx_result = await adapter.project(docx_path)

        assert "template_style_inventory" in dotx_result.metadata
        assert isinstance(dotx_result.metadata["template_style_inventory"], list)
        assert len(dotx_result.metadata["template_style_inventory"]) > 0

        assert "template_style_inventory" not in docx_result.metadata

    async def test_ad_071_inventory_entries_have_required_shape(self, tmp_path):
        """AD-071: Every inventory entry carries exactly the six required fields."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        dotx_path = _build_template_fixture(tmp_path, "shape.dotx")

        result = await adapter.project(dotx_path)
        inventory = result.metadata["template_style_inventory"]

        required_keys = {
            "id",
            "name",
            "type",
            "based_on",
            "has_numbering",
            "is_custom",
            "numbering_detail",
        }
        allowed_types = {"paragraph", "character", "table", "numbering"}

        for entry in inventory:
            assert isinstance(entry, dict)
            assert set(entry.keys()) == required_keys, (
                f"entry keys {set(entry.keys())} != required {required_keys}"
            )
            assert isinstance(entry["id"], str) and entry["id"]
            assert isinstance(entry["name"], str) and entry["name"]
            assert entry["type"] in allowed_types
            assert entry["based_on"] is None or isinstance(entry["based_on"], str)
            assert isinstance(entry["has_numbering"], bool)
            assert isinstance(entry["is_custom"], bool)
            # numbering_detail internal consistency: dict iff has_numbering
            if entry["has_numbering"]:
                assert isinstance(entry["numbering_detail"], dict)
            else:
                assert entry["numbering_detail"] is None

        # Verify the fixture produced a mix of custom and built-in
        customs = [e for e in inventory if e["is_custom"]]
        builtins = [e for e in inventory if not e["is_custom"]]
        assert customs, "fixture should contain at least one custom style"
        assert builtins, "fixture should contain at least one built-in style"

        # Verify the AnnotationChar style is recognized as character type,
        # not paragraph. Regression guard against the str(enum).split(".")
        # type-detection bug.
        ann = next((e for e in inventory if e["id"] == "AnnotationChar"), None)
        assert ann is not None, "AnnotationChar custom style missing from inventory"
        assert ann["type"] == "character"
        assert ann["based_on"] == "Normal"

    async def test_ad_072_has_numbering_reflects_active_reference(self, tmp_path):
        """AD-072: has_numbering is True only for styles with an active numPr."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        dotx_path = _build_template_fixture(tmp_path, "numbering.dotx")

        result = await adapter.project(dotx_path)
        inventory = result.metadata["template_style_inventory"]
        by_id = {e["id"]: e for e in inventory}
        by_name = {e["name"]: e for e in inventory}

        assert by_id["DefinitionsHeader"]["has_numbering"] is True, (
            "custom paragraph with active numPr should have has_numbering=True"
        )
        assert by_id["AppendixHeading"]["has_numbering"] is False, (
            "custom paragraph with no numPr should have has_numbering=False"
        )
        assert by_id["AnnotationChar"]["has_numbering"] is False, (
            "character style should always have has_numbering=False"
        )
        # Built-in Heading 1 has XML id "Heading1" but human name "Heading 1"
        assert by_name["Heading 1"]["has_numbering"] is True, (
            "built-in Heading 1 wired to active numbering should be True"
        )

    async def test_ad_073_dotx_emits_tags_docx_does_not(self, tmp_path):
        """AD-073:.dotx emits adapter_tags with prescribed namespacing;.docx does not."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        dotx_path = _build_template_fixture(tmp_path, "tags.dotx")

        # Equivalent.docx (custom styles added but no template content type)
        doc = docx.Document()
        doc.add_paragraph("Intro", style="Heading 1")
        _add_custom_style(doc, "AppendixHeading", "Appendix Heading", "paragraph")
        docx_path = tmp_path / "tags.docx"
        doc.save(str(docx_path))

        dotx_result = await adapter.project(dotx_path)
        docx_result = await adapter.project(docx_path)

        tags = dotx_result.metadata["adapter_tags"]
        assert isinstance(tags, list)

        # Style-presence tags: custom styles only
        assert "template:style:Appendix Heading" in tags
        assert "template:style:Definitions Header" in tags
        assert "template:style:Annotation Char" in tags
        # Built-in styles must not appear in style-presence namespace
        assert "template:style:Heading 1" not in tags
        assert "template:style:Normal" not in tags

        # Numbering tags: any style with has_numbering=True
        assert "template:has_numbering:Definitions Header" in tags
        assert "template:has_numbering:Heading 1" in tags  # built-in wired
        assert "template:has_numbering:Appendix Heading" not in tags

        # Adapter declares its owned prefixes so ingestion can strip stale
        # adapter-owned tags on force re-ingest.
        assert dotx_result.metadata["adapter_tag_prefixes"] == ["template:"]

        # .docx does not emit adapter_tags
        assert "adapter_tags" not in docx_result.metadata
        assert "adapter_tag_prefixes" not in docx_result.metadata

    async def test_ad_074_dotx_synthesizes_non_empty_text(self, tmp_path):
        """AD-074:.dotx projection produces non-empty style-surface text."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()
        dotx_path = _build_template_fixture(tmp_path, "surface.dotx")

        result = await adapter.project(dotx_path)

        text = result.text
        assert text.strip(), "template projection must produce non-empty text"

        # First line identifies the artifact as a Word template and includes
        # the title (filename stem here; the helper does not inject a
        # Title-styled paragraph).
        first_line = text.splitlines()[0]
        assert "template" in first_line.lower()
        assert "surface" in first_line.lower() or "Surface" in first_line

        # Every custom style's human name appears in the text
        inventory = result.metadata["template_style_inventory"]
        custom_names = [e["name"] for e in inventory if e["is_custom"]]
        for name in custom_names:
            assert name in text, f"custom style {name!r} missing from synthesized text"

        # Built-in Heading 1 wiring is called out
        assert "auto-numbering" in text or "auto-numbered" in text
        assert "Heading 1" in text

        # Built-in unmodified styles (Normal, Header, Footer) are NOT listed
        # in the style-name sections -- they would flood every template's
        # text with identical keyword noise. (The text MAY mention them in
        # passing, e.g., via a basedOn chain, but the style-list sections
        # should not enumerate them.)
        styles_section_end = text.find("Built-in styles")
        styles_section = text[:styles_section_end] if styles_section_end >= 0 else text
        assert "Normal" not in styles_section
        assert ": Header," not in styles_section  # avoid matching "Header Char"
        assert ": Footer," not in styles_section

    async def test_ad_075_numbering_detail_resolved_from_abstract(self, tmp_path):
        """AD-075: numbering_detail carries num_id, abstract_num_id, ilvl,
        num_fmt, lvl_text, and suppressed sourced from the resolved
        numbering definition."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        # Build a fixture with distinct numbering definitions exercising
        # decimal at ilvl 0, upperRoman at ilvl 1, and a lvlOverride that
        # suppresses numbering (numFmt=none).
        doc = docx.Document()
        doc.add_paragraph("Placeholder", style="Heading 1")

        abstract_xml = """
        <w:abstractNum w:abstractNumId="300"
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
          <w:lvl w:ilvl="0">
            <w:start w:val="1"/>
            <w:numFmt w:val="decimal"/>
            <w:lvlText w:val="%1."/>
            <w:lvlJc w:val="left"/>
          </w:lvl>
          <w:lvl w:ilvl="1">
            <w:start w:val="1"/>
            <w:numFmt w:val="upperRoman"/>
            <w:lvlText w:val="%1.%2"/>
            <w:lvlJc w:val="left"/>
          </w:lvl>
        </w:abstractNum>
        """
        num_decimal = """
        <w:num w:numId="300"
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
          <w:abstractNumId w:val="300"/>
        </w:num>
        """
        num_roman = """
        <w:num w:numId="301"
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
          <w:abstractNumId w:val="300"/>
        </w:num>
        """
        num_suppressed = """
        <w:num w:numId="302"
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
          <w:abstractNumId w:val="300"/>
          <w:lvlOverride w:ilvl="0">
            <w:lvl w:ilvl="0">
              <w:start w:val="1"/>
              <w:numFmt w:val="none"/>
              <w:lvlText w:val=""/>
              <w:lvlJc w:val="left"/>
            </w:lvl>
          </w:lvlOverride>
        </w:num>
        """
        # Inject these into a fresh numbering part
        doc.add_paragraph("dummy", style="List Bullet")
        doc.element.body.remove(doc.element.body[-1])
        numbering_elm = doc.part.numbering_part.numbering_definitions._numbering
        for xml in (abstract_xml, num_decimal, num_roman, num_suppressed):
            numbering_elm.append(etree.fromstring(xml))

        _add_custom_style(doc, "DecimalHeading", "Decimal Heading", "paragraph", num_id=300)
        _add_custom_style(doc, "RomanSubheading", "Roman Subheading", "paragraph", num_id=301)
        _add_custom_style(doc, "SuppressedHeading", "Suppressed Heading", "paragraph", num_id=302)
        _add_custom_style(doc, "PlainBody", "Plain Body", "paragraph")

        # The default num_id=300 builds with ilvl=0 via _add_custom_style;
        # override RomanSubheading's ilvl to 1 by direct XML edit.
        roman_style = doc.styles["Roman Subheading"]
        numPr = roman_style.element.find(qn("w:pPr")).find(qn("w:numPr"))
        numPr.find(qn("w:ilvl")).set(qn("w:val"), "1")

        docx_tmp = tmp_path / "numdetail_source.docx"
        dotx_path = tmp_path / "numdetail.dotx"
        doc.save(str(docx_tmp))
        _convert_docx_to_dotx(docx_tmp, dotx_path)
        docx_tmp.unlink()

        result = await DocxAdapter().project(dotx_path)
        inventory = result.metadata["template_style_inventory"]
        by_id = {e["id"]: e for e in inventory}

        # Decimal heading
        dec = by_id["DecimalHeading"]
        assert dec["has_numbering"] is True
        detail = dec["numbering_detail"]
        assert set(detail.keys()) == {
            "num_id",
            "abstract_num_id",
            "ilvl",
            "num_fmt",
            "lvl_text",
            "suppressed",
        }
        assert detail["num_id"] == 300
        assert detail["abstract_num_id"] == 300
        assert detail["ilvl"] == 0
        assert detail["num_fmt"] == "decimal"
        assert detail["lvl_text"] == "%1."
        assert detail["suppressed"] is False

        # Roman subheading (same abstract, different ilvl)
        roman = by_id["RomanSubheading"]
        rdetail = roman["numbering_detail"]
        assert rdetail["num_id"] == 301
        assert rdetail["abstract_num_id"] == 300
        assert rdetail["ilvl"] == 1
        assert rdetail["num_fmt"] == "upperRoman"
        assert rdetail["lvl_text"] == "%1.%2"
        assert rdetail["suppressed"] is False

        # Suppressed heading: same abstract, lvlOverride forces numFmt=none
        sup = by_id["SuppressedHeading"]
        assert sup["has_numbering"] is True, (
            "style still references an active num; has_numbering remains True"
        )
        sdetail = sup["numbering_detail"]
        assert sdetail["num_id"] == 302
        assert sdetail["num_fmt"] == "none"
        assert sdetail["suppressed"] is True

        # Plain body has no numbering
        plain = by_id["PlainBody"]
        assert plain["has_numbering"] is False
        assert plain["numbering_detail"] is None

    async def test_ad_074_docx_projection_unchanged(self, tmp_path):
        """AD-074 (negative):.docx projection does not receive synthesized text."""
        from sage.source_adapters.docx_adapter import DocxAdapter

        adapter = DocxAdapter()

        doc = docx.Document()
        doc.add_paragraph("Real Heading", style="Heading 1")
        doc.add_paragraph("Body paragraph content.")
        docx_path = tmp_path / "plain.docx"
        doc.save(str(docx_path))

        result = await adapter.project(docx_path)

        # Text contains only the body content, not template synthesis markers
        assert "Real Heading" in result.text
        assert "Body paragraph content." in result.text
        assert "Microsoft Word template" not in result.text
        assert "user-authored" not in result.text


# ── XLSX Adapter ──────────────────────────────────────────────────

try:
    import openpyxl

    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False

requires_openpyxl = pytest.mark.skipif(not _HAS_OPENPYXL, reason="openpyxl not available")


def _make_xlsx(tmp_path: Path, filename: str = "test.xlsx") -> Path:
    """Create a minimal single-sheet workbook and return its path."""
    wb = openpyxl.Workbook()
    path = tmp_path / filename
    wb.save(str(path))
    wb.close()
    return path


def _make_multisheet_xlsx(
    tmp_path: Path,
    sheets_data: dict[str, list[list]],
    filename: str = "multi.xlsx",
) -> Path:
    """Create a workbook with named sheets and cell data.

    sheets_data: {"SheetName": [[row1_values], [row2_values],...]}
    """
    wb = openpyxl.Workbook()
    # Remove default sheet
    wb.remove(wb.active)
    for sheet_name, rows in sheets_data.items():
        ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)
    path = tmp_path / filename
    wb.save(str(path))
    wb.close()
    return path


@requires_openpyxl
class TestXlsxAdapter:
    """AD-053 through AD-067: XLSX source adapter tests."""

    async def test_ad_053_basic_projection(self, tmp_path):
        """AD-053: Basic projection returns valid ProjectionResult."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {"Sales": [["Product", "Revenue"], ["Widget", 100], ["Gadget", 200]]},
            filename="basic.xlsx",
        )

        result = await adapter.project(path)

        assert isinstance(result.text, str)
        assert len(result.text) > 0
        assert len(result.headings) == 1
        assert result.headings[0].text == "Sales"
        assert result.headings[0].level == 1
        assert isinstance(result.content_hash, str)
        assert len(result.content_hash) == 64  # SHA-256 hex
        assert result.adapter_version == XlsxAdapter.VERSION

    async def test_ad_054_multisheet_headings(self, tmp_path):
        """AD-054: Multiple sheets produce one level-1 HeadingNode each."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {
                "Revenue": [["Q1", "Q2"], [100, 200]],
                "Expenses": [["Q1", "Q2"], [50, 80]],
                "Summary": [["Net"], [170]],
            },
        )

        result = await adapter.project(path)

        assert len(result.headings) == 3
        assert [h.text for h in result.headings] == ["Revenue", "Expenses", "Summary"]
        assert all(h.level == 1 for h in result.headings)

    async def test_ad_055_heading_paths_flat(self, tmp_path):
        """AD-055: Heading paths are sheet names only (no hierarchy)."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {
                "Alpha": [["A"], [1]],
                "Beta": [["B"], [2]],
            },
        )

        result = await adapter.project(path)

        assert result.headings[0].path == "Alpha"
        assert result.headings[1].path == "Beta"
        assert " > " not in result.headings[0].path
        assert " > " not in result.headings[1].path

    async def test_ad_056_column_headers_in_content(self, tmp_path):
        """AD-056: First row rendered as pipe-delimited header row."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {"Data": [["Name", "Age", "City"], ["Alice", 30, "NYC"]]},
            filename="headers.xlsx",
        )

        result = await adapter.project(path)

        assert "| Name | Age | City |" in result.headings[0].content

    async def test_ad_057_preview_rows_default(self, tmp_path):
        """AD-057: Default config includes first 5 data rows, omits row 6+."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        rows = [["ID", "Value"]] + [[i, f"val_{i}"] for i in range(1, 9)]
        path = _make_multisheet_xlsx(tmp_path, {"Data": rows}, filename="preview.xlsx")

        result = await adapter.project(path)
        content = result.headings[0].content

        # Rows 1-5 present
        for i in range(1, 6):
            assert f"val_{i}" in content
        # Rows 6-8 omitted
        for i in range(6, 9):
            assert f"val_{i}" not in content

    async def test_ad_058_preview_rows_config(self, tmp_path):
        """AD-058: preview_rows config limits data rows per sheet."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        rows = [["ID", "Value"]] + [[i, f"val_{i}"] for i in range(1, 9)]
        path = _make_multisheet_xlsx(tmp_path, {"Data": rows}, filename="limited.xlsx")

        result = await adapter.project(path, config={"preview_rows": 2})
        content = result.headings[0].content

        assert "val_1" in content
        assert "val_2" in content
        assert "val_3" not in content

    async def test_ad_059_dimensions_in_content(self, tmp_path):
        """AD-059: Content includes dimensions line for each sheet."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        rows = [["A", "B", "C"]] + [[1, 2, 3]] * 20
        path = _make_multisheet_xlsx(tmp_path, {"Big": rows}, filename="dims.xlsx")

        result = await adapter.project(path)
        content = result.headings[0].content

        assert "21 rows" in content
        assert "3 columns" in content

    async def test_ad_060_title_from_first_sheet(self, tmp_path):
        """AD-060: Title extracted from first sheet name."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {"Quarterly Report": [["Q1"], [100]], "Details": [["X"], [1]]},
            filename="report.xlsx",
        )

        result = await adapter.project(path)
        assert result.title == "Quarterly Report"

    async def test_ad_061_title_fallback_to_filename(self, tmp_path):
        """AD-061: Default sheet name falls back to filename stem."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        # openpyxl default sheet name is "Sheet"
        path = _make_xlsx(tmp_path, filename="my_data.xlsx")

        result = await adapter.project(path)
        assert result.title == "my_data"

    async def test_ad_062_content_hash_raw_bytes(self, tmp_path):
        """AD-062: content_hash is SHA-256 of raw.xlsx file bytes."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {"Data": [["X"], [1]]},
            filename="hash_test.xlsx",
        )

        expected_hash = hashlib.sha256(path.read_bytes()).hexdigest()
        result = await adapter.project(path)
        assert result.content_hash == expected_hash

    async def test_ad_063_source_modified_at(self, tmp_path):
        """AD-063: source_modified_at extracted from file mtime, timezone-aware."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {"Data": [["X"], [1]]},
            filename="provenance.xlsx",
        )

        known_mtime = datetime(2023, 9, 1, 14, 0, 0, tzinfo=timezone.utc)
        os.utime(path, (path.stat().st_atime, known_mtime.timestamp()))

        result = await adapter.project(path)

        assert "source_modified_at" in result.metadata
        parsed = datetime.fromisoformat(result.metadata["source_modified_at"])
        assert parsed.tzinfo is not None
        assert abs((parsed - known_mtime).total_seconds()) < 1.0

    async def test_ad_064_empty_workbook(self, tmp_path):
        """AD-064: Empty workbook produces valid result with filename-stem title."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_xlsx(tmp_path, filename="empty_book.xlsx")

        result = await adapter.project(path)

        assert result.title == "empty_book"
        assert isinstance(result.text, str)
        assert isinstance(result.content_hash, str)
        assert len(result.content_hash) == 64

    async def test_ad_065_max_sheets_config(self, tmp_path):
        """AD-065: max_sheets config limits number of sheets projected."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        sheets = {f"Sheet_{i}": [["Col"], [i]] for i in range(1, 6)}
        path = _make_multisheet_xlsx(tmp_path, sheets, filename="many.xlsx")

        result = await adapter.project(path, config={"max_sheets": 2})

        assert len(result.headings) == 2
        assert result.headings[0].text == "Sheet_1"
        assert result.headings[1].text == "Sheet_2"

    async def test_ad_066_sheet_metadata(self, tmp_path):
        """AD-066: metadata includes sheet_names, total_sheets, dimensions."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {
                "Alpha": [["A", "B"], [1, 2], [3, 4]],
                "Beta": [["X"], [10]],
            },
        )

        result = await adapter.project(path)

        assert result.metadata["sheet_names"] == ["Alpha", "Beta"]
        assert result.metadata["total_sheets"] == 2
        assert "Alpha" in result.metadata["dimensions"]
        assert "Beta" in result.metadata["dimensions"]

    async def test_ad_067_full_text_concatenation(self, tmp_path):
        """AD-067: result.text concatenates all sheet projections with markdown headings."""
        from sage.source_adapters.xlsx_adapter import XlsxAdapter

        adapter = XlsxAdapter()
        path = _make_multisheet_xlsx(
            tmp_path,
            {
                "Revenue": [["Q1", "Q2"], [100, 200]],
                "Costs": [["Q1", "Q2"], [50, 80]],
            },
        )

        result = await adapter.project(path)

        # Sheet names appear as markdown headings in full text
        assert "# Revenue" in result.text
        assert "# Costs" in result.text
        # Content from both sheets present
        assert "100" in result.text
        assert "50" in result.text


# ── PDF Adapter ───────────────────────────────────────────────────

try:
    import pdfplumber as _pdfplumber  # noqa: F401
    import pypdf as _pypdf  # noqa: F401

    _HAS_PDFPLUMBER = True
except ImportError:
    _HAS_PDFPLUMBER = False

try:
    from reportlab.pdfgen import canvas as _rl_canvas  # noqa: F401

    _HAS_REPORTLAB = True
except ImportError:
    _HAS_REPORTLAB = False

try:
    from PIL import Image as _PilImage  # noqa: F401

    _HAS_PIL = True
except ImportError:
    _HAS_PIL = False

try:
    import ocrmypdf as _ocrmypdf  # noqa: F401

    _HAS_OCRMYPDF = True
except ImportError:
    _HAS_OCRMYPDF = False

import shutil as _shutil_for_ocr  # noqa: E402 -- grouped with PDF section

_HAS_OCR_BINARIES = bool(_shutil_for_ocr.which("tesseract") and _shutil_for_ocr.which("gs"))


requires_pdf = pytest.mark.skipif(
    not (_HAS_PDFPLUMBER and _HAS_REPORTLAB),
    reason="pdfplumber, pypdf, or reportlab not available",
)
requires_pdf_with_image = pytest.mark.skipif(
    not (_HAS_PDFPLUMBER and _HAS_REPORTLAB and _HAS_PIL),
    reason="pdfplumber, pypdf, reportlab, or Pillow not available",
)
requires_ocr = pytest.mark.skipif(
    not (_HAS_OCRMYPDF and _HAS_OCR_BINARIES),
    reason="ocrmypdf or tesseract/ghostscript not available",
)


def _draw_paragraph_lines(c, lines: list[str], start_y: int = 750) -> None:
    """Draw lines of text on the current page, top-down."""
    y = start_y
    for line in lines:
        c.drawString(72, y, line)
        y -= 14


def _make_pdf_with_pages(
    path: Path,
    pages: list[list[str]],
    info_title: str | None = None,
) -> Path:
    """Create a multi-page PDF with given lines per page.

    pages: [[line1, line2,...], [line1, line2,...],...]
    Each inner list is the lines of one page.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    # reportlab defaults /Info /Title to "untitled" if not set; explicitly
    # clear it when the caller wants no /Info /Title.
    c.setTitle(info_title if info_title is not None else "")
    for page_lines in pages:
        _draw_paragraph_lines(c, page_lines)
        c.showPage()
    c.save()
    return path


def _make_pdf_with_outline(
    path: Path,
    outline: list[tuple[int, str, int]],
    pages: list[list[str]],
    info_title: str | None = None,
) -> Path:
    """Create a multi-page PDF with a bookmark outline.

    outline: list of (level_1_indexed, title, page_index_0_indexed).
    pages: text lines per page.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    # reportlab defaults /Info /Title to "untitled" if not set; explicitly
    # clear it when the caller wants no /Info /Title.
    c.setTitle(info_title if info_title is not None else "")

    bookmarks_per_page: dict[int, list[tuple[int, str, str]]] = {}
    for idx, (level, title, page_idx) in enumerate(outline):
        key = f"BM_{idx}"
        bookmarks_per_page.setdefault(page_idx, []).append((level, title, key))

    parent_keys: list[str] = []
    parent_levels: list[int] = []

    for page_idx, page_lines in enumerate(pages):
        for level, title, key in bookmarks_per_page.get(page_idx, []):
            c.bookmarkPage(key)
            while parent_levels and parent_levels[-1] >= level:
                parent_keys.pop()
                parent_levels.pop()
            c.addOutlineEntry(title, key, level=level - 1, closed=False)
            parent_keys.append(key)
            parent_levels.append(level)
        _draw_paragraph_lines(c, page_lines)
        c.showPage()

    c.showOutline()
    c.save()
    return path


def _make_scanned_pdf(path: Path, image_caption: str = "scanned text") -> Path:
    """Create a single-page PDF whose only content is a rasterized image.

    No text drawing operations are performed, so pdfplumber extracts an
    empty (or near-empty) text layer.
    """
    from PIL import Image, ImageDraw
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    img_path = path.with_suffix(".png")
    img = Image.new("RGB", (400, 100), color="white")
    draw = ImageDraw.Draw(img)
    draw.text((10, 40), image_caption, fill="black")
    img.save(img_path)

    c = canvas.Canvas(str(path), pagesize=letter)
    c.setTitle("")
    c.drawImage(str(img_path), 72, 600, width=400, height=100)
    c.showPage()
    c.save()
    img_path.unlink()
    return path


def _make_encrypted_pdf(path: Path, body_text: str = "encrypted content") -> Path:
    """Create a PDF and encrypt it with an owner password."""
    import pypdf
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    plain_path = path.with_suffix(".plain.pdf")
    c = canvas.Canvas(str(plain_path), pagesize=letter)
    c.drawString(72, 720, body_text)
    c.showPage()
    c.save()

    reader = pypdf.PdfReader(str(plain_path))
    writer = pypdf.PdfWriter(clone_from=reader)
    writer.encrypt(user_password="userpw", owner_password="ownerpw")
    with open(path, "wb") as f:
        writer.write(f)
    plain_path.unlink()
    return path


def _make_zero_page_pdf(path: Path) -> Path:
    """Create a structurally-valid PDF with zero pages."""
    import pypdf

    writer = pypdf.PdfWriter()
    with open(path, "wb") as f:
        writer.write(f)
    return path


def _make_corrupt_pdf(path: Path) -> Path:
    """Create a file that begins with a PDF header but is not a valid PDF."""
    path.write_bytes(b"%PDF-1.7\nthis is not a valid PDF body content\n")
    return path


def _make_malformed_xref_pdf(path: Path) -> Path:
    """Create a minimal PDF with intentionally-wrong xref offsets.

    pypdf reads such PDFs by falling back to an object scan; the recovery
    succeeds but emits stderr warnings ("incorrect startxref pointer",
    "parsing for Object Streams", or similar). The adapter is required to
    suppress these.
    """
    pdf_bytes = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]"
        b" /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
        b"4 0 obj\n<< /Length 55 >>\nstream\n"
        b"BT /F1 12 Tf 50 700 Td (MALFORMED_XREF_BODY) Tj ET\n"
        b"endstream\nendobj\n"
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000999 00000 n \n"
        b"0000000999 00000 n \n"
        b"0000000999 00000 n \n"
        b"0000000999 00000 n \n"
        b"0000000999 00000 n \n"
        b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
        b"startxref\n500\n%%EOF\n"
    )
    path.write_bytes(pdf_bytes)
    return path


# Unit tests for the PDF adapter's CID safe-decode helper.


def test_decode_safe_cid_decodes_ascii_range():
    from sage.source_adapters.pdf_adapter import _decode_safe_cid

    cid_form = (
        "(cid:77)(cid:65)(cid:76)(cid:70)(cid:79)(cid:82)(cid:77)"
        "(cid:69)(cid:68)(cid:95)(cid:88)(cid:82)(cid:69)(cid:70)"
        "(cid:95)(cid:66)(cid:79)(cid:68)(cid:89)"
    )
    assert _decode_safe_cid(cid_form) == "MALFORMED_XREF_BODY"


def test_decode_safe_cid_preserves_non_ascii_cids():
    from sage.source_adapters.pdf_adapter import _decode_safe_cid

    assert _decode_safe_cid("(cid:200)(cid:300)") == "(cid:200)(cid:300)"


def test_decode_safe_cid_preserves_control_range_cids():
    from sage.source_adapters.pdf_adapter import _decode_safe_cid

    assert _decode_safe_cid("(cid:0)(cid:7)(cid:31)") == "(cid:0)(cid:7)(cid:31)"


def test_decode_safe_cid_passthrough_unicode():
    from sage.source_adapters.pdf_adapter import _decode_safe_cid

    assert _decode_safe_cid("Hello world") == "Hello world"


def test_decode_safe_cid_mixed_safe_and_unsafe():
    from sage.source_adapters.pdf_adapter import _decode_safe_cid

    assert _decode_safe_cid("(cid:65)hello(cid:200)world(cid:66)") == "Ahello(cid:200)worldB"


@requires_pdf
class TestPdfAdapter:
    """AD-076 through AD-094: PDF source adapter tests (v0.1, native-text only)."""

    # ── Section 8.1 — Registration & basic projection ────────────

    def test_ad_076_extension_registration(self):
        """AD-076: PdfAdapter registers.pdf as a supported extension."""
        import re

        from sage.source_adapters.pdf_adapter import PdfAdapter

        assert PdfAdapter.EXTENSIONS == [".pdf"]
        assert isinstance(PdfAdapter.VERSION, str)
        assert re.match(r"^\d+\.\d+\.\d+$", PdfAdapter.VERSION)

    async def test_ad_077_native_text_flat_projection(self, tmp_path):
        """AD-077: Native-text PDF without outline produces single flat heading."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_pages(
            tmp_path / "flat.pdf",
            pages=[["UNIQUE_BODY_MARKER hello there", "second line"]],
        )
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert "UNIQUE_BODY_MARKER" in result.text
        assert len(result.headings) == 1
        assert result.headings[0].level == 1
        assert result.headings[0].path == result.title
        assert "UNIQUE_BODY_MARKER" in result.headings[0].content

    async def test_ad_078_content_hash_is_sha256_of_bytes(self, tmp_path):
        """AD-078: content_hash is SHA-256 hex of raw source bytes."""
        import hashlib

        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_pages(tmp_path / "hash.pdf", pages=[["hash test body"]])
        adapter = PdfAdapter()
        result = await adapter.project(path)

        expected = hashlib.sha256(path.read_bytes()).hexdigest()
        assert result.content_hash == expected
        assert len(result.content_hash) == 64
        assert result.content_hash == result.content_hash.lower()

    # ── Section 8.2 — Provenance & metadata ──────────────────────

    async def test_ad_079_source_modified_at_from_mtime(self, tmp_path):
        """AD-079: source_modified_at extracted from file mtime."""
        import os
        from datetime import datetime, timezone

        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_pages(tmp_path / "mtime.pdf", pages=[["mtime test"]])
        # Set a known mtime
        target_ts = datetime(2025, 6, 15, 12, 30, 0, tzinfo=timezone.utc).timestamp()
        os.utime(path, (target_ts, target_ts))

        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert isinstance(result.metadata["source_modified_at"], str)
        parsed = datetime.fromisoformat(result.metadata["source_modified_at"])
        assert parsed.tzinfo is not None
        assert (
            abs((parsed - datetime.fromtimestamp(target_ts, tz=timezone.utc)).total_seconds()) < 1
        )

    async def test_ad_080_page_count_metadata(self, tmp_path):
        """AD-080: page_count metadata reflects actual page count."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_pages(
            tmp_path / "multipage.pdf",
            pages=[[f"page {i} body"] for i in range(7)],
        )
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert result.metadata["page_count"] == 7
        assert result.metadata["pages_extracted"] == 7
        assert "pdf:truncated" not in result.metadata.get("adapter_tags", [])

    # ── Section 8.3 — Title priority chain ───────────────────────

    async def test_ad_081_title_from_info_title(self, tmp_path):
        """AD-081: Title resolves from /Info /Title when present (priority 1)."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_outline(
            tmp_path / "info_title.pdf",
            outline=[(1, "Outline Heading 1", 0)],
            pages=[["First Line Title", "more body"]],
            info_title="Set Via Info",
        )
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert result.title == "Set Via Info"

    async def test_ad_082_title_from_first_outline_entry(self, tmp_path):
        """AD-082: Title falls back to first outline entry when /Info absent."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_outline(
            tmp_path / "outline_title.pdf",
            outline=[(1, "Outline Heading 1", 0)],
            pages=[["First Line Title", "more body"]],
            info_title=None,
        )
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert result.title == "Outline Heading 1"

    async def test_ad_083_title_from_first_body_line(self, tmp_path):
        """AD-083: Title falls back to first body line <=120 chars."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_pages(
            tmp_path / "body_title.pdf",
            pages=[["First Line Title", "more content here"]],
            info_title=None,
        )
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert result.title == "First Line Title"

    async def test_ad_084_title_falls_back_to_filename(self, tmp_path):
        """AD-084: Title falls back to filename stem when no usable line."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        long_line = "X" * 200
        path = _make_pdf_with_pages(
            tmp_path / "long-leading-line.pdf",
            pages=[[long_line]],
            info_title=None,
        )
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert result.title == "long-leading-line"

    # ── Section 8.4 — Outline handling ───────────────────────────

    async def test_ad_085_outline_produces_nested_headings(self, tmp_path):
        """AD-085: Outlined PDF produces HeadingNodes mirroring outline structure."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_outline(
            tmp_path / "outlined.pdf",
            outline=[
                (1, "Intro", 0),
                (2, "Background", 1),
                (3, "Prior Art", 2),
            ],
            pages=[
                ["INTRO_BODY"],
                ["BACKGROUND_BODY"],
                ["PRIOR_ART_BODY"],
            ],
            info_title="Outlined Doc",
        )
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert len(result.headings) == 3
        assert result.headings[0].level == 1
        assert result.headings[0].text == "Intro"
        assert result.headings[0].path == "Intro"
        assert "INTRO_BODY" in result.headings[0].content

        assert result.headings[1].level == 2
        assert result.headings[1].path == "Intro > Background"
        assert "BACKGROUND_BODY" in result.headings[1].content

        assert result.headings[2].level == 3
        assert result.headings[2].path == "Intro > Background > Prior Art"
        assert "PRIOR_ART_BODY" in result.headings[2].content

    async def test_ad_086_has_outline_metadata_and_tag(self, tmp_path):
        """AD-086: has_outline=True and pdf:has_outline tag emitted with outline."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        outlined = _make_pdf_with_outline(
            tmp_path / "with_outline.pdf",
            outline=[(1, "Section 1", 0)],
            pages=[["body"]],
            info_title="With Outline",
        )
        no_outline = _make_pdf_with_pages(
            tmp_path / "no_outline.pdf",
            pages=[["body"]],
            info_title="No Outline",
        )
        adapter = PdfAdapter()

        outlined_result = await adapter.project(outlined)
        no_outline_result = await adapter.project(no_outline)

        assert outlined_result.metadata["has_outline"] is True
        assert "pdf:has_outline" in outlined_result.metadata.get("adapter_tags", [])
        assert no_outline_result.metadata["has_outline"] is False
        assert "pdf:has_outline" not in no_outline_result.metadata.get("adapter_tags", [])

    async def test_ad_087_outline_depth_cap_drops_deeper_entries(self, tmp_path):
        """AD-087: Outline depth cap (10) drops deeper entries; text preserved in ancestor."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        outline = [(level, f"Level{level}", level - 1) for level in range(1, 13)]
        pages = [[f"BODY_L{level}"] for level in range(1, 13)]

        path = _make_pdf_with_outline(
            tmp_path / "deep.pdf",
            outline=outline,
            pages=pages,
            info_title="Deep Outline",
        )
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert max(node.level for node in result.headings) == 10
        heading_texts = [node.text for node in result.headings]
        assert "Level11" not in heading_texts
        assert "Level12" not in heading_texts

        level10_node = next(n for n in result.headings if n.level == 10)
        assert "BODY_L11" in level10_node.content
        assert "BODY_L12" in level10_node.content

    # ── Section 8.5 — Scanned-PDF OCR pre-pass ────────────────────

    @requires_pdf_with_image
    @requires_ocr
    async def test_ad_088_scanned_pdf_produces_text_with_ocr_applied_tag(self, tmp_path):
        """AD-088: Scanned PDF gets inline OCR; projection has text and pdf:ocr_applied."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_scanned_pdf(tmp_path / "scanned-doc.pdf")
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert result.text != "", "OCR pre-pass must produce non-empty text"
        adapter_tags = result.metadata.get("adapter_tags", [])
        assert "pdf:ocr_applied" in adapter_tags
        assert "pdf:scanned" not in adapter_tags
        assert "pdf:ocr_no_text" not in adapter_tags

    @requires_pdf_with_image
    @requires_ocr
    async def test_ad_089_adapter_tag_prefixes_declared(self, tmp_path):
        """AD-089: adapter_tag_prefixes declares ['pdf:'] for the OCR'd path too."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        outlined = _make_pdf_with_outline(
            tmp_path / "outlined_for_prefixes.pdf",
            outline=[(1, "Section 1", 0)],
            pages=[["body"]],
            info_title="With Outline",
        )
        scanned = _make_scanned_pdf(tmp_path / "scanned_for_prefixes.pdf")
        adapter = PdfAdapter()

        outlined_result = await adapter.project(outlined)
        scanned_result = await adapter.project(scanned)

        assert outlined_result.metadata.get("adapter_tag_prefixes") == ["pdf:"]
        assert scanned_result.metadata.get("adapter_tag_prefixes") == ["pdf:"]

    @requires_pdf_with_image
    async def test_ad_ocr_offloaded_to_worker_thread(self, tmp_path, monkeypatch):
        """OCR runs on the dedicated ``sage-ocr`` executor thread, not the event
        loop. The blocking pre-pass must be dispatched via run_in_executor like
        the embedding/abstraction stages, so a large scan cannot stall the loop
        (and the container liveness probe) mid-ingest. Monkeypatching the OCR
        step lets this run without the real toolchain."""
        import threading

        from sage.source_adapters import pdf_adapter as pdf_adapter_mod
        from sage.source_adapters.pdf_adapter import PdfAdapter

        recorded: dict[str, threading.Thread] = {}

        def _spy_ocr(source_path):
            # Record the thread the OCR pre-pass runs on, then hand back a PDF
            # carrying known text so the post-OCR re-extraction yields a real
            # pdf:ocr_applied projection (proving the offload is transparent).
            recorded["thread"] = threading.current_thread()
            out = tmp_path / "ocr-output.pdf"
            _make_pdf_with_pages(out, [["OCRDONE marker text"]])
            return out

        monkeypatch.setattr(pdf_adapter_mod, "_ocr_to_tempfile", _spy_ocr)

        path = _make_scanned_pdf(tmp_path / "scanned-offload.pdf")
        result = await PdfAdapter().project(path)

        worker = recorded.get("thread")
        assert worker is not None, "the OCR pre-pass never ran"
        assert worker is not threading.main_thread(), (
            "OCR ran on the main thread (still on the loop)"
        )
        assert worker.name.startswith("sage-ocr"), (
            f"OCR ran off the dedicated executor: thread={worker.name!r}"
        )
        # Offload is transparent to projection output.
        assert "pdf:ocr_applied" in result.metadata.get("adapter_tags", [])
        assert "OCRDONE" in result.text

    @requires_pdf_with_image
    async def test_ad_ocr_bounds_passed_to_ocrmypdf(self, tmp_path, monkeypatch):
        """The OCR call is bounded for the cor-prod container (2 vCPU / 4 GiB
        alongside the resident embedder): ``ocrmypdf.ocr`` receives an explicit
        ``jobs`` worker cap and a per-page ``tesseract_timeout``. A future edit
        that drops the bounds would let a large scan fan out unbounded."""
        import sys
        import types
        from pathlib import Path

        from sage.source_adapters import pdf_adapter as pdf_adapter_mod
        from sage.source_adapters.pdf_adapter import PdfAdapter

        captured: dict[str, object] = {}

        def _capturing_ocr(input_path, output_path, *args, **kwargs):
            captured.update(kwargs)
            Path(output_path).write_bytes(b"%PDF-1.4\n%%EOF\n")

        fake_ocrmypdf = types.ModuleType("ocrmypdf")
        fake_ocrmypdf.ocr = _capturing_ocr
        monkeypatch.setitem(sys.modules, "ocrmypdf", fake_ocrmypdf)

        # First extraction (against the source) trips is_scanned; the second
        # (against the OCR output) returns empty so we do not need a real OCR'd
        # PDF -- the assertion is on the kwargs the OCR call received.
        original_extract = pdf_adapter_mod._extract_from_path
        calls = {"n": 0}

        def _stub_extract(p, max_pages):
            calls["n"] += 1
            return original_extract(p, max_pages) if calls["n"] == 1 else ([""], [], None, 1, 1)

        monkeypatch.setattr(pdf_adapter_mod, "_extract_from_path", _stub_extract)

        path = _make_scanned_pdf(tmp_path / "scanned-bounds.pdf")
        await PdfAdapter().project(path)

        assert "jobs" in captured, "ocrmypdf.ocr must receive an explicit jobs worker cap"
        assert captured["jobs"] == 1, f"jobs not bounded to 1: {captured.get('jobs')!r}"
        assert "tesseract_timeout" in captured, "ocrmypdf.ocr must receive a per-page timeout"
        assert isinstance(captured["tesseract_timeout"], (int, float))
        assert captured["tesseract_timeout"] > 0

    @requires_pdf_with_image
    async def test_ad_095_scanned_pdf_without_ocrmypdf_raises(self, tmp_path, monkeypatch):
        """AD-095: Scanned PDF with ocrmypdf unimportable raises ValueError naming [ocr] extra."""
        import sys

        from sage.source_adapters.pdf_adapter import PdfAdapter

        # Force the lazy `import ocrmypdf` inside the OCR helper to fail.
        monkeypatch.setitem(sys.modules, "ocrmypdf", None)

        path = _make_scanned_pdf(tmp_path / "scanned-noocr.pdf")
        adapter = PdfAdapter()

        with pytest.raises(ValueError) as exc_info:
            await adapter.project(path)

        message = str(exc_info.value)
        assert "[ocr]" in message
        assert "tesseract" in message
        assert "ghostscript" in message

    @requires_pdf_with_image
    async def test_ad_096_ocrmypdf_runtime_error_raises_and_cleans_tempfile(
        self, tmp_path, monkeypatch
    ):
        """AD-096: ocrmypdf.ocr raising mid-call → ValueError; tempfile is unlinked."""
        import sys
        import tempfile as _tempfile
        import types
        from pathlib import Path

        from sage.source_adapters import pdf_adapter as pdf_adapter_mod
        from sage.source_adapters.pdf_adapter import PdfAdapter

        # Fake ocrmypdf module whose.ocr() raises mid-call.
        fake_ocrmypdf = types.ModuleType("ocrmypdf")

        def _failing_ocr(*args, **kwargs):
            raise RuntimeError("simulated tesseract failure")

        fake_ocrmypdf.ocr = _failing_ocr
        monkeypatch.setitem(sys.modules, "ocrmypdf", fake_ocrmypdf)

        # Capture the tempfile path the helper creates so we can assert
        # it has been unlinked after the raise propagates.
        captured_paths: list[str] = []
        real_named_tempfile = _tempfile.NamedTemporaryFile

        def _spy_named_tempfile(*args, **kwargs):
            handle = real_named_tempfile(*args, **kwargs)
            captured_paths.append(handle.name)
            return handle

        monkeypatch.setattr(pdf_adapter_mod.tempfile, "NamedTemporaryFile", _spy_named_tempfile)

        path = _make_scanned_pdf(tmp_path / "scanned-runtime-error.pdf")
        adapter = PdfAdapter()

        with pytest.raises(ValueError, match="OCR failed for"):
            await adapter.project(path)

        assert captured_paths, "OCR helper must allocate a tempfile before invoking ocrmypdf"
        for p in captured_paths:
            assert not Path(p).exists(), f"tempfile {p} leaked after OCR failure"

    @requires_pdf_with_image
    async def test_ad_097_ocr_yields_no_text_emits_ocr_no_text_tag(self, tmp_path, monkeypatch):
        """AD-097: OCR succeeds but yields no extractable text → pdf:ocr_no_text tag."""
        import sys
        import types

        from sage.source_adapters import pdf_adapter as pdf_adapter_mod
        from sage.source_adapters.pdf_adapter import PdfAdapter

        # Stub ocrmypdf so the OCR call is a no-op (succeeds without touching the file).
        fake_ocrmypdf = types.ModuleType("ocrmypdf")

        def _noop_ocr(input_path, output_path, *args, **kwargs):
            # Make the output path a valid (if empty) PDF so any unlink in
            # the success path succeeds; content is irrelevant because we
            # also stub the post-OCR extraction.
            from pathlib import Path

            Path(output_path).write_bytes(b"%PDF-1.4\n%%EOF\n")

        fake_ocrmypdf.ocr = _noop_ocr
        monkeypatch.setitem(sys.modules, "ocrmypdf", fake_ocrmypdf)

        # Wrap _extract_from_path: first call (against the source) runs the
        # real extraction so is_scanned trips; second call (against the
        # OCR tempfile) returns empty page texts to simulate a blank scan.
        original_extract = pdf_adapter_mod._extract_from_path
        call_count = {"n": 0}

        def _stub_extract(path, max_pages):
            call_count["n"] += 1
            if call_count["n"] == 1:
                return original_extract(path, max_pages)
            return ([""], [], None, 1, 1)

        monkeypatch.setattr(pdf_adapter_mod, "_extract_from_path", _stub_extract)

        path = _make_scanned_pdf(tmp_path / "scanned-blank.pdf")
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert call_count["n"] == 2, "OCR path must invoke _extract_from_path twice"
        assert result.text == ""
        assert result.headings == []
        adapter_tags = result.metadata.get("adapter_tags", [])
        assert "pdf:ocr_no_text" in adapter_tags
        assert "pdf:ocr_applied" not in adapter_tags
        assert "pdf:scanned" not in adapter_tags

    @pytest.mark.parametrize(
        "tempdir, env_tmpdir, expected_kind",
        [
            ("/tmp", "/tmp/x", "cache"),  # (A) under /tmp, $TMPDIR under /tmp -> ~/.cache
            ("/tmp", "/var/folders/safe", "env"),  # (B) under /tmp, $TMPDIR off /tmp -> $TMPDIR
            ("/var/folders/zz", "/tmp/x", "none"),  # (C) already off /tmp -> no-op (None)
        ],
    )
    async def test_ad_107_safe_ocr_tempdir_selection(
        self, monkeypatch, tempdir, env_tmpdir, expected_kind
    ):
        """AD-107: _safe_ocr_tempdir picks a non-/tmp base only when the process
        temp dir resolves under /tmp; otherwise it is a no-op."""
        import os
        import tempfile

        from sage.source_adapters import pdf_adapter as pdf_adapter_mod

        # Keep the selector hermetic: never touch the real filesystem.
        monkeypatch.setattr(pdf_adapter_mod.os, "makedirs", lambda *a, **k: None)
        monkeypatch.setattr(tempfile, "tempdir", tempdir, raising=False)
        monkeypatch.setenv("TMPDIR", env_tmpdir)

        result = pdf_adapter_mod._safe_ocr_tempdir()

        if expected_kind == "cache":
            assert result == os.path.expanduser("~/.cache")
            assert not result.startswith(("/tmp", "/private/tmp"))
        elif expected_kind == "env":
            assert result == env_tmpdir
            assert not result.startswith(("/tmp", "/private/tmp"))
        else:
            assert result is None

    async def test_ad_108_ocr_intermediate_raster_routed_off_tmp(self, tmp_path, monkeypatch):
        """AD-108: _ocr_to_tempfile routes the OCR intermediate raster and the
        output tempfile off /tmp when the process temp dir resolves under /tmp
        (leptonica cannot read a /tmp-rooted raster on macOS)."""
        import sys
        import tempfile
        import types
        from pathlib import Path

        from sage.source_adapters import pdf_adapter as pdf_adapter_mod

        # Force the failure precondition: process temp dir under /tmp.
        monkeypatch.setattr(tempfile, "tempdir", "/tmp", raising=False)
        monkeypatch.setenv("TMPDIR", "/tmp/cas-regress")

        # Positive control: prove the precondition is real, so a pass below
        # reflects the adapter's rerouting, not a coincidentally-safe env.
        assert tempfile.gettempdir().startswith("/tmp")

        # Fake ocrmypdf records the effective tempdir at the moment ocr() runs.
        seen: dict[str, str] = {}
        fake_ocrmypdf = types.ModuleType("ocrmypdf")

        def _spy_ocr(input_path, output_path, *args, **kwargs):
            seen["tempdir"] = tempfile.gettempdir()
            Path(output_path).write_bytes(b"%PDF-1.4\n%%EOF\n")

        fake_ocrmypdf.ocr = _spy_ocr
        monkeypatch.setitem(sys.modules, "ocrmypdf", fake_ocrmypdf)

        src = tmp_path / "scanned.pdf"
        src.write_bytes(b"%PDF-1.4 dummy")

        out = pdf_adapter_mod._ocr_to_tempfile(src)
        try:
            assert "tempdir" in seen, "ocrmypdf.ocr must have been invoked"
            assert not seen["tempdir"].startswith(("/tmp", "/private/tmp")), (
                "OCR intermediates must be routed off /tmp"
            )
            assert not str(out).startswith(("/tmp", "/private/tmp")), (
                "OCR output tempfile must be routed off /tmp"
            )
        finally:
            out.unlink(missing_ok=True)

    # ── Section 8.6 — Failure modes ──────────────────────────────

    async def test_ad_090_encrypted_pdf_raises(self, tmp_path):
        """AD-090: Encrypted PDF raises ValueError."""
        import re

        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_encrypted_pdf(tmp_path / "encrypted.pdf")
        adapter = PdfAdapter()
        with pytest.raises(ValueError, match=re.compile(r"encrypted", re.IGNORECASE)):
            await adapter.project(path)

    async def test_ad_091_corrupt_pdf_raises(self, tmp_path):
        """AD-091: Corrupt PDF raises ValueError."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        bad_header = _make_corrupt_pdf(tmp_path / "garbage.pdf")
        empty = tmp_path / "empty.pdf"
        empty.write_bytes(b"")

        adapter = PdfAdapter()

        with pytest.raises(ValueError):
            await adapter.project(bad_header)
        with pytest.raises(ValueError):
            await adapter.project(empty)

    async def test_ad_092_malformed_pdf_no_stderr_leakage(self, tmp_path, caplog):
        """AD-092: Malformed-but-readable PDF projects without log/stderr leakage.

        Asserts via ``caplog`` rather than ``contextlib.redirect_stderr``:
        in production, pypdf/pdfminer emit malformed-xref warnings through
        ``logging`` (not direct ``sys.stderr`` writes); the adapter's
        suppression sets those loggers to ERROR so WARN records never
        propagate to any handler. ``caplog`` intercepts records before
        any handler, so it correctly observes whether the level filter
        is doing its job. Under pytest's logging plugin, the records
        would never have reached ``sys.stderr`` anyway, which made the
        old ``redirect_stderr``-based assertion non-load-bearing.
        """
        import logging

        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_malformed_xref_pdf(tmp_path / "malformed.pdf")
        adapter = PdfAdapter()

        with (
            caplog.at_level(logging.WARNING, logger="pypdf"),
            caplog.at_level(logging.WARNING, logger="pdfminer"),
        ):
            result = await adapter.project(path)

        assert "MALFORMED_XREF_BODY" in result.text
        log_messages = " ".join(record.getMessage() for record in caplog.records)
        assert "incorrect startxref" not in log_messages
        assert "Ignoring wrong pointing object" not in log_messages
        assert "parsing for Object Streams" not in log_messages

    async def test_ad_093_zero_page_pdf_empty_projection(self, tmp_path):
        """AD-093: Empty (zero-page) PDF produces empty projection without error."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_zero_page_pdf(tmp_path / "zero-pages.pdf")
        adapter = PdfAdapter()
        result = await adapter.project(path)

        assert result.text == ""
        assert result.headings == []
        assert result.metadata["page_count"] == 0
        assert result.metadata["pages_extracted"] == 0
        assert result.title == "zero-pages"

    # ── Section 8.7 — Configuration ──────────────────────────────

    async def test_ad_094_max_pages_truncation(self, tmp_path):
        """AD-094: max_pages truncates; pages_extracted and pdf:truncated tag
        distinguish source vs. projection."""
        from sage.source_adapters.pdf_adapter import PdfAdapter

        path = _make_pdf_with_pages(
            tmp_path / "ten_pages.pdf",
            pages=[[f"PAGE_{i}_BODY"] for i in range(1, 11)],
        )
        adapter = PdfAdapter()

        truncated = await adapter.project(path, config={"max_pages": 3})
        no_truncation = await adapter.project(path, config={"max_pages": 10})
        above_actual = await adapter.project(path, config={"max_pages": 100})

        # Truncation case
        assert truncated.metadata["page_count"] == 10
        assert truncated.metadata["pages_extracted"] == 3
        assert "pdf:truncated" in truncated.metadata.get("adapter_tags", [])
        assert "PAGE_1_BODY" in truncated.text
        assert "PAGE_2_BODY" in truncated.text
        assert "PAGE_3_BODY" in truncated.text
        assert "PAGE_4_BODY" not in truncated.text
        assert "PAGE_10_BODY" not in truncated.text

        # No-truncation case
        assert no_truncation.metadata["page_count"] == 10
        assert no_truncation.metadata["pages_extracted"] == 10
        assert "pdf:truncated" not in no_truncation.metadata.get("adapter_tags", [])

        # Limit-above-actual case
        assert above_actual.metadata["page_count"] == 10
        assert above_actual.metadata["pages_extracted"] == 10
        assert "pdf:truncated" not in above_actual.metadata.get("adapter_tags", [])


# ── PPTX Adapter ──────────────────────────────────────────────────

try:
    import pptx as _pptx  # noqa: F401

    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False


requires_pptx = pytest.mark.skipif(not _HAS_PPTX, reason="python-pptx not available")
requires_pptx_with_image = pytest.mark.skipif(
    not (_HAS_PPTX and _HAS_PIL), reason="python-pptx or Pillow not available"
)

# Layout indices in python-pptx's default template.
_LAYOUT_TITLE_ONLY = 5
_LAYOUT_BLANK = 6

# OPC main-part content types for the presentation and template flavors.
_PPTX_MAIN_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)
_POTX_MAIN_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"


def _new_presentation():
    from pptx import Presentation

    return Presentation()


def _add_slide(prs, title: str | None):
    """Add a slide, using the blank layout when no title is wanted."""
    layout = prs.slide_layouts[_LAYOUT_TITLE_ONLY if title is not None else _LAYOUT_BLANK]
    slide = prs.slides.add_slide(layout)
    if title is not None:
        slide.shapes.title.text = title
    return slide


def _add_textbox(slide, text: str, top_in: float = 2.0, left_in: float = 1.0):
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(3), Inches(1))
    box.text_frame.text = text
    return box


def _make_pptx(tmp_path, slides: list[dict], filename: str = "test.pptx") -> Path:
    """Build a .pptx from a compact slide spec.

    Each entry of ``slides`` is a dict with optional keys:

    - ``title``: title-placeholder text; omit or pass None for a blank
      layout with no title placeholder at all.
    - ``body``: list of strings, each rendered as its own text box,
      stacked top-down so authoring order matches visual order.
    - ``notes``: speaker-notes text.
    - ``table``: list of rows (each a list of cell strings).
    """
    prs = _new_presentation()
    for spec in slides:
        slide = _add_slide(prs, spec.get("title"))
        top = 2.0
        for line in spec.get("body", []):
            _add_textbox(slide, line, top_in=top)
            top += 1.0
        if "table" in spec:
            _add_pptx_table(slide, spec["table"], top_in=top)
        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]
    path = tmp_path / filename
    prs.save(str(path))
    return path


def _add_pptx_table(slide, rows: list[list[str]], top_in: float = 4.0):
    from pptx.util import Inches

    frame = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(1), Inches(top_in), Inches(5), Inches(1)
    )
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            frame.table.cell(r, c).text = cell
    return frame


def _make_empty_pptx(tmp_path, filename: str = "empty.pptx") -> Path:
    """A presentation carrying zero slides."""
    prs = _new_presentation()
    path = tmp_path / filename
    prs.save(str(path))
    return path


def _make_potx(tmp_path, slides: list[dict], filename: str = "template.potx") -> Path:
    """Build a .pptx then rewrite its main content type to the template flavor."""
    import zipfile as _zipfile

    source = _make_pptx(tmp_path, slides, filename="__potx_source.pptx")
    path = tmp_path / filename
    with _zipfile.ZipFile(source, "r") as z_in:
        with _zipfile.ZipFile(path, "w", _zipfile.ZIP_DEFLATED) as z_out:
            for item in z_in.namelist():
                data = z_in.read(item)
                if item == "[Content_Types].xml":
                    data = data.replace(
                        _PPTX_MAIN_TYPE.encode("utf-8"), _POTX_MAIN_TYPE.encode("utf-8")
                    )
                z_out.writestr(item, data)
    source.unlink()
    return path


def _make_corrupt_pptx(path: Path) -> Path:
    """ZIP magic bytes followed by garbage: valid header, unreadable package."""
    path.write_bytes(b"PK\x03\x04this is not a valid OPC package at all\n")
    return path


def _make_encrypted_pptx(path: Path) -> Path:
    """OLE2 compound-file magic, which is how password-protected OOXML is stored."""
    path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 64)
    return path


def _strip_position(shape) -> None:
    """Remove a shape's explicit geometry so ``top``/``left`` report None."""
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
    xfrm = shape._element.spPr.find(ns)
    shape._element.spPr.remove(xfrm)


def _set_alt_text(shape, text: str) -> None:
    shape._element._nvXxPr.cNvPr.set("descr", text)


def _slide_heading(result, n: int):
    """Return the level-1 heading whose text starts with ``Slide {n}``."""
    for heading in result.headings:
        if heading.level == 1 and heading.text.startswith(f"Slide {n}"):
            return heading
    raise AssertionError(f"no level-1 heading for slide {n} in {[h.text for h in result.headings]}")


@requires_pptx
class TestPptxAdapter:
    """AD-098 through AD-128: PPTX source adapter tests."""

    # ── Section 11.1 — Registration and projection shape ──────────

    async def test_ad_098_extensions_and_version(self):
        """AD-098: Adapter declares its extensions and a semver VERSION."""
        import re

        from sage.source_adapters.pptx_adapter import PptxAdapter

        assert PptxAdapter.EXTENSIONS == [".pptx", ".potx"]
        assert isinstance(PptxAdapter.VERSION, str)
        assert re.match(r"^\d+\.\d+\.\d+$", PptxAdapter.VERSION)

    async def test_ad_099_basic_projection(self, tmp_path):
        """AD-099: Basic projection returns a valid ProjectionResult."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [
                {"title": "Alpha", "body": ["BODY_ONE"]},
                {"title": "Beta", "body": ["BODY_TWO"]},
                {"title": "Gamma", "body": ["BODY_THREE"]},
            ],
            filename="basic.pptx",
        )

        result = await PptxAdapter().project(path)

        assert isinstance(result.text, str)
        assert len(result.text) > 0
        assert len([h for h in result.headings if h.level == 1]) == 3
        assert isinstance(result.content_hash, str)
        assert len(result.content_hash) == 64  # SHA-256 hex
        assert result.adapter_version == PptxAdapter.VERSION
        assert result.metadata["slide_count"] == 3
        assert result.metadata["slides_projected"] == 3

    # ── Section 11.2 — Heading synthesis ──────────────────────────

    async def test_ad_100_one_heading_per_slide_in_deck_order(self, tmp_path):
        """AD-100: One level-1 heading per slide, numbered, in deck order."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"title": "Alpha"}, {"title": "Beta"}, {"title": "Gamma"}],
            filename="ordered.pptx",
        )

        result = await PptxAdapter().project(path)
        top_level = [h for h in result.headings if h.level == 1]

        assert [h.text for h in top_level] == [
            "Slide 1: Alpha",
            "Slide 2: Beta",
            "Slide 3: Gamma",
        ]
        for heading in top_level:
            assert heading.path == heading.text
            assert " > " not in heading.path
        # The title placeholder is consumed by the heading, so a title-only
        # slide projects an empty body rather than repeating its own title.
        assert [h.content for h in top_level] == ["", "", ""]

    async def test_ad_101_untitled_slide_does_not_disturb_numbering(self, tmp_path):
        """AD-101: A slide with no title placeholder does not shift its neighbours' numbers.

        Originally this asserted a bare ``Slide 2`` heading for the untitled
        slide. Title inference (AD-124) intentionally changed that: an untitled
        slide carrying title-shaped body text now takes its heading from that
        text. The bare-placeholder path is still asserted, on the inputs where
        inference correctly declines -- AD-102 (no text at all), AD-125 (line
        too long), AD-126 (a table). What this case still owns is the numbering
        invariant, which no other test covers with a titled/untitled mix.
        """
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"title": "Alpha"}, {"body": ["ORPHAN_BODY"]}, {"title": "Gamma"}],
            filename="untitled.pptx",
        )

        result = await PptxAdapter().project(path)
        top_level = [h for h in result.headings if h.level == 1]

        assert [h.text for h in top_level] == [
            "Slide 1: Alpha",
            "Slide 2: ORPHAN_BODY",
            "Slide 3: Gamma",
        ]

    async def test_ad_102_empty_slide_still_projects_its_heading(self, tmp_path):
        """AD-102: A slide with no shapes keeps its heading so numbering stays contiguous."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"title": "Alpha", "body": ["BODY_ONE"]}, {}, {"title": "Gamma"}],
            filename="gap.pptx",
        )

        result = await PptxAdapter().project(path)
        top_level = [h for h in result.headings if h.level == 1]

        assert len(top_level) == 3
        assert top_level[1].text == "Slide 2"
        assert top_level[1].content == ""
        # Numbering after the gap must not shift.
        assert top_level[2].text == "Slide 3: Gamma"

    # ── Section 11.3 — Speaker notes ──────────────────────────────

    async def test_ad_103_notes_project_under_a_subheading(self, tmp_path):
        """AD-103: Speaker notes project as a level-2 heading addressable via its path."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [
                {"title": "Alpha"},
                {"title": "Beta", "body": ["BODY_TWO"], "notes": "NOTES_FOR_BETA"},
            ],
            filename="notes.pptx",
        )

        result = await PptxAdapter().project(path)
        notes = [h for h in result.headings if h.text == "Notes"]

        assert len(notes) == 1
        assert notes[0].level == 2
        assert notes[0].path == "Slide 2: Beta > Notes"
        assert notes[0].content == "NOTES_FOR_BETA"
        # The slide body stays on the slide heading, not the notes heading.
        assert "BODY_TWO" in _slide_heading(result, 2).content
        assert "BODY_TWO" not in notes[0].content
        assert result.metadata["notes_count"] == 1

    async def test_ad_104_no_notes_emits_no_subheading(self, tmp_path):
        """AD-104: A deck with no speaker notes emits no Notes heading."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"title": "Alpha", "body": ["BODY_ONE"]}, {"title": "Beta"}],
            filename="nonotes.pptx",
        )

        result = await PptxAdapter().project(path)

        assert [h for h in result.headings if h.text == "Notes"] == []
        assert result.metadata["notes_count"] == 0

    async def test_ad_105_has_notes_tag_tracks_notes_presence(self, tmp_path):
        """AD-105: pptx:has_notes is emitted iff at least one slide carries notes."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        adapter = PptxAdapter()
        with_notes = _make_pptx(
            tmp_path, [{"title": "Alpha", "notes": "SOME_NOTES"}], filename="withnotes.pptx"
        )
        without_notes = _make_pptx(tmp_path, [{"title": "Alpha"}], filename="withoutnotes.pptx")

        assert "pptx:has_notes" in (await adapter.project(with_notes)).metadata["adapter_tags"]
        assert (
            "pptx:has_notes" not in (await adapter.project(without_notes)).metadata["adapter_tags"]
        )

    # ── Section 11.4 — Content recovery ───────────────────────────

    async def test_ad_106_table_projects_as_markdown(self, tmp_path):
        """AD-106: Tables project as Markdown tables."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"title": "Alpha", "table": [["Name", "Role"], ["Ada", "Engineer"]]}],
            filename="table.pptx",
        )

        content = _slide_heading(await PptxAdapter().project(path), 1).content

        assert "| Name | Role |" in content
        assert "| --- | --- |" in content
        assert "| Ada | Engineer |" in content

    async def test_ad_107_group_shape_text_is_recovered(self, tmp_path):
        """AD-107: Text inside a grouped shape is recovered, not dropped."""
        from pptx.util import Inches

        from sage.source_adapters.pptx_adapter import PptxAdapter

        prs = _new_presentation()
        slide = _add_slide(prs, "Alpha")
        group = slide.shapes.add_group_shape()
        member = group.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
        member.text_frame.text = "GROUPED_TEXT"
        path = tmp_path / "group.pptx"
        prs.save(str(path))

        content = _slide_heading(await PptxAdapter().project(path), 1).content

        assert "GROUPED_TEXT" in content

    async def test_ad_108_nested_group_text_is_recovered(self, tmp_path):
        """AD-108: Text inside a group nested within a group is recovered."""
        from pptx.util import Inches

        from sage.source_adapters.pptx_adapter import PptxAdapter

        prs = _new_presentation()
        slide = _add_slide(prs, "Alpha")
        outer = slide.shapes.add_group_shape()
        inner = outer.shapes.add_group_shape()
        member = inner.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
        member.text_frame.text = "DEEPLY_NESTED_TEXT"
        path = tmp_path / "nested.pptx"
        prs.save(str(path))

        content = _slide_heading(await PptxAdapter().project(path), 1).content

        assert "DEEPLY_NESTED_TEXT" in content

    @requires_pptx_with_image
    async def test_ad_109_alt_text_contributes_but_filenames_do_not(self, tmp_path):
        """AD-109: Authored alt text contributes; a default filename descr does not."""
        import io

        from PIL import Image
        from pptx.util import Inches

        from sage.source_adapters.pptx_adapter import PptxAdapter

        prs = _new_presentation()
        slide = _add_slide(prs, "Alpha")
        buf = io.BytesIO()
        Image.new("RGB", (20, 20), "white").save(buf, format="PNG")

        buf.seek(0)
        described = slide.shapes.add_picture(buf, Inches(1), Inches(2), Inches(1), Inches(1))
        _set_alt_text(described, "AUTHORED_ALT_TEXT")

        buf.seek(0)
        undescribed = slide.shapes.add_picture(buf, Inches(1), Inches(4), Inches(1), Inches(1))
        _set_alt_text(undescribed, "diagram_v2.png")

        path = tmp_path / "altext.pptx"
        prs.save(str(path))

        content = _slide_heading(await PptxAdapter().project(path), 1).content

        assert "AUTHORED_ALT_TEXT" in content
        # PowerPoint defaults descr to the image filename; that is not alt text.
        assert "diagram_v2.png" not in content

    # ── Section 11.5 — Reading order ──────────────────────────────

    async def test_ad_110_shapes_read_top_down_then_left_right(self, tmp_path):
        """AD-110: Shapes project in visual reading order, not shape-tree order."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        prs = _new_presentation()
        slide = _add_slide(prs, "Alpha")
        # Authored bottom-first, and right-before-left within the shared row,
        # so shape-tree order is the exact reverse of reading order.
        _add_textbox(slide, "ZULU_BOTTOM", top_in=5.0, left_in=1.0)
        _add_textbox(slide, "MIKE_ROW_RIGHT", top_in=3.0, left_in=6.0)
        _add_textbox(slide, "MIKE_ROW_LEFT", top_in=3.05, left_in=1.0)
        _add_textbox(slide, "ALFA_TOP", top_in=1.5, left_in=1.0)
        path = tmp_path / "order.pptx"
        prs.save(str(path))

        content = _slide_heading(await PptxAdapter().project(path), 1).content
        positions = [
            content.index(marker)
            for marker in ("ALFA_TOP", "MIKE_ROW_LEFT", "MIKE_ROW_RIGHT", "ZULU_BOTTOM")
        ]

        assert positions == sorted(positions), content

    async def test_ad_111_unresolved_position_sorts_last_in_tree_order(self, tmp_path):
        """AD-111: Shapes with no resolvable position sort last, keeping tree order."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        prs = _new_presentation()
        slide = _add_slide(prs, "Alpha")
        floating_first = _add_textbox(slide, "FLOATING_ONE", top_in=1.0)
        floating_second = _add_textbox(slide, "FLOATING_TWO", top_in=1.5)
        _add_textbox(slide, "POSITIONED_BODY", top_in=5.0)
        _strip_position(floating_first)
        _strip_position(floating_second)
        path = tmp_path / "unpositioned.pptx"
        prs.save(str(path))

        content = _slide_heading(await PptxAdapter().project(path), 1).content

        # Positioned shape wins despite sitting lowest on the slide...
        assert content.index("POSITIONED_BODY") < content.index("FLOATING_ONE")
        # ...and the unpositioned pair keeps its shape-tree order.
        assert content.index("FLOATING_ONE") < content.index("FLOATING_TWO")

    # ── Section 11.6 — Provenance and title ───────────────────────

    async def test_ad_112_content_hash_is_sha256_of_source_bytes(self, tmp_path):
        """AD-112: content_hash is the SHA-256 of the raw source bytes."""
        import hashlib

        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(tmp_path, [{"title": "Alpha", "body": ["BODY"]}], filename="hash.pptx")

        result = await PptxAdapter().project(path)

        assert result.content_hash == hashlib.sha256(path.read_bytes()).hexdigest()

    async def test_ad_113_source_modified_at_from_mtime(self, tmp_path):
        """AD-113: source_modified_at is the file mtime as a timezone-aware ISO string."""
        from datetime import datetime, timezone

        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(tmp_path, [{"title": "Alpha"}], filename="mtime.pptx")
        expected = datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc)

        result = await PptxAdapter().project(path)
        parsed = datetime.fromisoformat(result.metadata["source_modified_at"])

        assert parsed == expected
        assert parsed.tzinfo is not None

    async def test_ad_114_title_prefers_first_slide_title(self, tmp_path):
        """AD-114: Title comes from the first slide's title, else the filename stem."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        adapter = PptxAdapter()
        titled = _make_pptx(
            tmp_path, [{"title": "Deck Headline"}, {"title": "Second"}], filename="titled.pptx"
        )
        untitled = _make_pptx(tmp_path, [{"body": ["BODY"]}], filename="fallback-name.pptx")

        assert (await adapter.project(titled)).title == "Deck Headline"
        assert (await adapter.project(untitled)).title == "fallback-name"

    # ── Section 11.7 — Adapter tags ───────────────────────────────

    async def test_ad_115_tag_prefix_declared_unconditionally(self, tmp_path):
        """AD-115: adapter_tag_prefixes is declared even when no tags are emitted."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path, [{"title": "Alpha", "body": ["BODY_ONE"]}], filename="notags.pptx"
        )

        result = await PptxAdapter().project(path)

        # A deck with text and no notes emits no tags at all...
        assert result.metadata["adapter_tags"] == []
        # ...but the namespace must still be declared, or a stale pptx: tag
        # from a prior ingest survives the re-ingest strip.
        assert result.metadata["adapter_tag_prefixes"] == ["pptx:"]

    async def test_ad_116_truncation_emits_tag_and_stops_projecting(self, tmp_path):
        """AD-116: max_slides truncates projection and emits pptx:truncated."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"title": f"Slide{i}", "body": [f"BODY_{i}"]} for i in range(1, 6)],
            filename="long.pptx",
        )
        adapter = PptxAdapter()

        truncated = await adapter.project(path, config={"max_slides": 2})
        full = await adapter.project(path, config={"max_slides": 5})

        assert len([h for h in truncated.headings if h.level == 1]) == 2
        assert truncated.metadata["slide_count"] == 5
        assert truncated.metadata["slides_projected"] == 2
        assert "pptx:truncated" in truncated.metadata["adapter_tags"]
        assert "BODY_2" in truncated.text
        assert "BODY_3" not in truncated.text

        assert "pptx:truncated" not in full.metadata["adapter_tags"]
        assert "BODY_5" in full.text

    @requires_pptx_with_image
    async def test_ad_117_no_text_tag_when_a_slide_yields_nothing(self, tmp_path):
        """AD-117: A slide with no recoverable text emits pptx:no_text and keeps its heading."""
        import io

        from PIL import Image
        from pptx.util import Inches

        from sage.source_adapters.pptx_adapter import PptxAdapter

        prs = _new_presentation()
        _add_slide(prs, "Alpha")
        bare = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_BLANK])
        buf = io.BytesIO()
        Image.new("RGB", (20, 20), "white").save(buf, format="PNG")
        buf.seek(0)
        picture = bare.shapes.add_picture(buf, Inches(1), Inches(2), Inches(1), Inches(1))
        _set_alt_text(picture, "screenshot.png")
        path = tmp_path / "imageonly.pptx"
        prs.save(str(path))

        result = await PptxAdapter().project(path)

        assert "pptx:no_text" in result.metadata["adapter_tags"]
        assert _slide_heading(result, 2).content == ""
        assert len([h for h in result.headings if h.level == 1]) == 2

    # ── Section 11.8 — Degenerate and error paths ─────────────────

    async def test_ad_118_zero_slide_deck_projects_empty(self, tmp_path):
        """AD-118: A deck with no slides projects empty rather than raising."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_empty_pptx(tmp_path, filename="no-slides.pptx")

        result = await PptxAdapter().project(path)

        assert result.text == ""
        assert result.headings == []
        assert result.title == "no-slides"
        assert result.metadata["slide_count"] == 0

    async def test_ad_119_potx_template_is_accepted(self, tmp_path):
        """AD-119: A .potx template opens via the content-type rewrite and projects."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_potx(
            tmp_path, [{"title": "Template Deck", "body": ["TEMPLATE_BODY"]}], filename="deck.potx"
        )

        result = await PptxAdapter().project(path)

        assert _slide_heading(result, 1).text == "Slide 1: Template Deck"
        assert "TEMPLATE_BODY" in result.text

    async def test_ad_120_corrupt_source_raises(self, tmp_path):
        """AD-120: A corrupt or empty deck raises ValueError with no partial output."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        garbage = _make_corrupt_pptx(tmp_path / "garbage.pptx")
        empty = tmp_path / "empty-bytes.pptx"
        empty.write_bytes(b"")
        adapter = PptxAdapter()

        with pytest.raises(ValueError):
            await adapter.project(garbage)
        with pytest.raises(ValueError):
            await adapter.project(empty)

    async def test_ad_121_locked_source_raises_naming_the_reason(self, tmp_path):
        """AD-121: A password-protected deck raises ValueError naming the reason.

        The source path is stripped from the message before matching. Every
        adapter interpolates the path into its diagnostics, and pytest names
        ``tmp_path`` after the test function, so a test whose own name
        contains the token being matched would pass on the directory name no
        matter what the adapter said.
        """
        import re

        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_encrypted_pptx(tmp_path / "locked.pptx")

        with pytest.raises(ValueError) as excinfo:
            await PptxAdapter().project(path)

        diagnostic = str(excinfo.value).replace(str(path), "")
        assert re.search(r"password|encrypt", diagnostic, re.IGNORECASE), diagnostic

    # ── Section 11.9 — Wiring ─────────────────────────────────────

    async def test_ad_122_pptx_is_a_binary_container_source(self):
        """AD-122: pptx is classified as a binary-container source type."""
        from sage.models.enums import BINARY_CONTAINER_SOURCE_TYPES, SourceType

        assert SourceType.PPTX in BINARY_CONTAINER_SOURCE_TYPES

    async def test_ad_123_adapter_is_registered_for_the_pptx_source_type(self):
        """AD-123: The pptx adapter is wired into the runtime adapter registry."""
        from sage.mcp_init import build_source_adapter_registry
        from sage.models.enums import SourceType
        from sage.source_adapters.pptx_adapter import PptxAdapter

        registry = build_source_adapter_registry()

        assert isinstance(registry[SourceType.PPTX], PptxAdapter)

    # ── Section 11.10 — Title inference for placeholder-free decks ──

    async def test_ad_124_untitled_slide_infers_title_from_first_body_line(self, tmp_path):
        """AD-124: A slide with no title placeholder takes its heading from the topmost text."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"body": ["The problem: knowledge is scattered", "Supporting detail."]}],
            filename="inferred.pptx",
        )

        result = await PptxAdapter().project(path)
        heading = _slide_heading(result, 1)

        assert heading.text == "Slide 1: The problem: knowledge is scattered"
        assert heading.path == heading.text
        # Promoted to the heading, so it must not also remain in the body.
        assert heading.content == "Supporting detail."

    async def test_ad_125_overlong_first_line_is_not_promoted(self, tmp_path):
        """AD-125: A first line too long to be a title leaves the heading as the placeholder."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        paragraph = (
            "This is a full sentence of body prose that runs on well past any "
            "reasonable title length and therefore must not be promoted into the "
            "slide heading where it would be unusable as an address."
        )
        assert len(paragraph) > 120
        path = _make_pptx(tmp_path, [{"body": [paragraph]}], filename="longline.pptx")

        result = await PptxAdapter().project(path)
        heading = _slide_heading(result, 1)

        assert heading.text == "Slide 1"
        # Not promoted, so it stays in the body rather than being lost.
        assert paragraph in heading.content

    async def test_ad_126_table_row_is_not_promoted_as_a_title(self, tmp_path):
        """AD-126: A slide whose topmost content is a table keeps the placeholder heading."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"table": [["Mode", "Returns"], ["semantic", "ranked hits"]]}],
            filename="tabletop.pptx",
        )

        result = await PptxAdapter().project(path)
        heading = _slide_heading(result, 1)

        assert heading.text == "Slide 1"
        assert "| Mode | Returns |" in heading.content

    async def test_ad_127_title_placeholder_wins_over_inference(self, tmp_path):
        """AD-127: A real title placeholder takes precedence; no body line is consumed."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path,
            [{"title": "Authored Title", "body": ["FIRST_BODY_LINE", "SECOND_BODY_LINE"]}],
            filename="authored.pptx",
        )

        result = await PptxAdapter().project(path)
        heading = _slide_heading(result, 1)

        assert heading.text == "Slide 1: Authored Title"
        # Inference must not fire, so the first body line survives in the body.
        assert heading.content == "FIRST_BODY_LINE\nSECOND_BODY_LINE"

    async def test_ad_128_document_title_still_falls_back_to_the_filename_stem(self, tmp_path):
        """AD-128: Heading inference does not change the document-title chain."""
        from sage.source_adapters.pptx_adapter import PptxAdapter

        path = _make_pptx(
            tmp_path, [{"body": ["Inferred Heading Text"]}], filename="stem-name.pptx"
        )

        result = await PptxAdapter().project(path)

        # The heading infers...
        assert _slide_heading(result, 1).text == "Slide 1: Inferred Heading Text"
        # ...but the document title stays on the placeholder-or-stem chain.
        assert result.title == "stem-name"
