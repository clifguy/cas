"""PPTX source adapter: extracts structured text from PowerPoint decks.

A deck has no intrinsic heading hierarchy the way a word-processing
document does, so the adapter synthesizes one: each slide becomes a
top-level heading carrying its slide number and title, with the slide's
body flattened beneath it and speaker notes nested one level below that.
Because the ingestion chunker splits on heading hierarchy, per-slide
headings yield per-slide chunk granularity, which matches how a deck is
actually cited ("slide 12 says...").

The slide number in the heading text is load-bearing rather than
decorative. ``HeadingNode.path`` is the address callers pass to
``read_section``, and slide titles repeat freely within a deck ("Agenda",
"Questions"), so a title-only path would collide. Numbering also stays
contiguous across slides that recover no text at all, so a caller can
always address slide N by its position in the deck.

Reading order is positional, not shape-tree order. A slide's shape tree
reflects z-order and authoring history: a text box added late sorts last
even when it sits at the top of the slide. Shapes are therefore sorted
top-to-bottom and then left-to-right, with a tolerance band so shapes on
one visual row read left-to-right rather than by sub-millimetre drift.
Shapes whose geometry resolves to nothing (a placeholder inheriting a
position the layout never sets) sort last; because the sort is stable,
those retain shape-tree order, so the algorithm degrades to shape-tree
order exactly where position is unknowable instead of scrambling.

Scope excludes rendering slides to images and offers no visual-fidelity
guarantee. This adapter recovers text and structure; a deck whose meaning
lives entirely in an unlabeled diagram projects thin, which is the
expected outcome rather than a defect.

Computes SHA-256 of raw source bytes for content_hash.
"""

import hashlib
import re
import shutil
import tempfile
import zipfile
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from sage.source_adapters.base import HeadingNode, ProjectionResult, SourceAdapter

# OPC main-part content types for the presentation and template flavors.
_PPTX_MAIN_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)
_POTX_MAIN_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"

# A password-protected OOXML file is an OLE2 compound document wrapping the
# encrypted package, not a ZIP, so it is identifiable by magic bytes alone.
_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
_ZIP_MAGIC = b"PK\x03\x04"

# PowerPoint populates a picture's `descr` with the source image's filename
# when the author never wrote alt text. That is a storage artifact, not a
# description, and projecting it would inject noise into every deck and mask
# genuinely text-free slides from the pptx:no_text signal.
_IMAGE_FILENAME_RE = re.compile(
    r"^[\w\-. ()]+\.(png|jpe?g|gif|bmp|tiff?|emf|wmf|svg|webp)$", re.IGNORECASE
)

# Vertical tolerance for treating two shapes as sharing a visual row: 0.25in
# in English Metric Units (914400 EMU per inch).
_ROW_TOLERANCE_EMU = 228600

_DEFAULT_MAX_SLIDES = 500

_ADAPTER_TAG_PREFIXES = ["pptx:"]


def _shape_position(shape) -> tuple[int, int] | None:
    """Return a shape's (top, left) in EMU, or None when unresolvable.

    Placeholders inherit geometry from their layout, which python-pptx
    resolves transparently; a shape carrying no explicit geometry and no
    inheritable source reports None.
    """
    try:
        top, left = shape.top, shape.left
    except (AttributeError, ValueError):
        return None
    if top is None:
        return None
    return (top, left if left is not None else 0)


def _reading_order_key(shape) -> tuple[int, int, int]:
    """Sort key placing shapes in visual reading order.

    Unresolvable positions sort into a trailing group; the caller's stable
    sort preserves shape-tree order within it.
    """
    position = _shape_position(shape)
    if position is None:
        return (1, 0, 0)
    top, left = position
    return (0, top // _ROW_TOLERANCE_EMU, left)


def _ordered_shapes(shapes) -> list:
    """Return shapes in visual reading order (stable within equal keys)."""
    return sorted(shapes, key=_reading_order_key)


def _alt_text(shape) -> str | None:
    """Return a shape's authored alt text, or None.

    python-pptx exposes no public accessor for the descriptive text on a
    shape's non-visual properties, so this reads the underlying element.
    Values that are merely the embedded image's filename are rejected.
    """
    try:
        descr = shape._element._nvXxPr.cNvPr.get("descr")
    except AttributeError:
        return None
    if not descr:
        return None
    descr = descr.strip()
    if not descr or _IMAGE_FILENAME_RE.match(descr):
        return None
    return descr


def _render_table(table) -> list[str]:
    """Render a table as Markdown rows, treating the first row as the header."""
    lines: list[str] = []
    for row_index, row in enumerate(table.rows):
        cells = [cell.text.replace("\n", " ").replace("|", r"\|").strip() for cell in row.cells]
        lines.append("| " + " | ".join(cells) + " |")
        if row_index == 0:
            lines.append("| " + " | ".join("---" for _ in cells) + " |")
    return lines


def _render_shape(shape) -> list[str]:
    """Flatten one shape to text lines, recursing into grouped shapes."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        lines: list[str] = []
        for member in _ordered_shapes(shape.shapes):
            lines.extend(_render_shape(member))
        return lines

    if shape.has_table:
        return _render_table(shape.table)

    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            return [text]

    alt = _alt_text(shape)
    return [alt] if alt else []


def _title_shape(slide):
    """Return the slide's title placeholder, or None when it has none."""
    try:
        return slide.shapes.title
    except (AttributeError, ValueError):
        return None


def _title_element(slide):
    """Return the underlying XML element of the slide's title placeholder.

    python-pptx constructs a fresh proxy object on every shape access, so a
    title placeholder can only be recognized during shape iteration by the
    element it wraps, never by object identity.
    """
    shape = _title_shape(slide)
    return None if shape is None else shape._element


def _slide_title(slide) -> str:
    """Return the slide's title text, or an empty string when it has none."""
    shape = _title_shape(slide)
    if shape is None or not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _slide_notes(slide) -> str:
    """Return the slide's speaker-notes text, or an empty string."""
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    if frame is None:
        return ""
    return frame.text.strip()


def _open_presentation(source_path: Path) -> Presentation:
    """Load a .pptx or .potx into a python-pptx Presentation.

    python-pptx rejects a template at load time because the main part's OPC
    content type is the template flavor rather than the presentation
    flavor. Templates are opened by copying the package into a temporary
    file with that one entry rewritten; every other part is untouched, so
    the rest of the adapter is unaffected.
    """
    try:
        with zipfile.ZipFile(source_path) as package:
            content_types = package.read("[Content_Types].xml")
    except (zipfile.BadZipFile, KeyError, OSError) as exc:
        raise ValueError(f"Failed to read presentation package {source_path}: {exc}") from exc

    if _POTX_MAIN_TYPE.encode("utf-8") not in content_types:
        try:
            return Presentation(str(source_path))
        except Exception as exc:
            raise ValueError(f"Failed to open presentation {source_path}: {exc}") from exc

    tmp_dir = Path(tempfile.mkdtemp(prefix="sage_potx_"))
    try:
        shadow = tmp_dir / "shadow.pptx"
        with zipfile.ZipFile(source_path) as z_in:
            with zipfile.ZipFile(shadow, "w", zipfile.ZIP_DEFLATED) as z_out:
                for item in z_in.namelist():
                    data = z_in.read(item)
                    if item == "[Content_Types].xml":
                        data = data.replace(
                            _POTX_MAIN_TYPE.encode("utf-8"), _PPTX_MAIN_TYPE.encode("utf-8")
                        )
                    z_out.writestr(item, data)
        try:
            return Presentation(str(shadow))
        except Exception as exc:
            raise ValueError(f"Failed to open presentation template {source_path}: {exc}") from exc
    finally:
        # python-pptx has read the package into memory by the time
        # Presentation() returns, so the temp dir is safe to remove.
        shutil.rmtree(tmp_dir, ignore_errors=True)


class PptxAdapter(SourceAdapter):
    VERSION = "0.1.0"
    EXTENSIONS = [".pptx", ".potx"]

    async def project(self, source_path: Path, config: dict | None = None) -> ProjectionResult:
        config = config or {}
        max_slides = config.get("max_slides", _DEFAULT_MAX_SLIDES)

        try:
            raw_bytes = source_path.read_bytes()
        except OSError as exc:
            raise ValueError(f"Failed to read presentation {source_path}: {exc}") from exc
        content_hash = hashlib.sha256(raw_bytes).hexdigest()

        if raw_bytes.startswith(_OLE2_MAGIC):
            raise ValueError(
                f"Presentation is password-protected or otherwise encrypted "
                f"and cannot be projected: {source_path}"
            )
        if not raw_bytes.startswith(_ZIP_MAGIC):
            raise ValueError(
                f"Presentation is not a readable OPC package (bad or truncated file): {source_path}"
            )

        presentation = _open_presentation(source_path)

        slides = list(presentation.slides)
        slide_count = len(slides)
        slides_to_project = slides[:max_slides] if max_slides else slides

        headings: list[HeadingNode] = []
        text_parts: list[str] = []
        notes_count = 0
        has_text_free_slide = False
        first_slide_title = ""

        for index, slide in enumerate(slides_to_project, start=1):
            title = _slide_title(slide)
            if index == 1:
                first_slide_title = title
            heading_text = f"Slide {index}: {title}" if title else f"Slide {index}"

            # Compared by underlying element, not by proxy identity: python-pptx
            # builds a fresh wrapper object on every shape access, so the title
            # fetched here is never the same Python object as the one yielded
            # during iteration even though both wrap the same shape.
            title_element = _title_element(slide)
            body_lines: list[str] = []
            for shape in _ordered_shapes(slide.shapes):
                if title_element is not None and shape._element is title_element:
                    continue
                body_lines.extend(_render_shape(shape))
            body = "\n".join(body_lines)

            if not body:
                has_text_free_slide = True

            headings.append(
                HeadingNode(level=1, text=heading_text, path=heading_text, content=body)
            )
            text_parts.append(f"# {heading_text}")
            if body:
                text_parts.append(body)

            notes = _slide_notes(slide)
            if notes:
                notes_count += 1
                headings.append(
                    HeadingNode(
                        level=2,
                        text="Notes",
                        path=f"{heading_text} > Notes",
                        content=notes,
                    )
                )
                text_parts.append("## Notes")
                text_parts.append(notes)

        adapter_tags: list[str] = []
        if notes_count:
            adapter_tags.append("pptx:has_notes")
        if has_text_free_slide:
            adapter_tags.append("pptx:no_text")
        if len(slides_to_project) < slide_count:
            adapter_tags.append("pptx:truncated")

        source_mtime = datetime.fromtimestamp(source_path.stat().st_mtime, tz=timezone.utc)

        return ProjectionResult(
            text="\n\n".join(text_parts),
            headings=headings,
            content_hash=content_hash,
            adapter_version=self.VERSION,
            title=first_slide_title or source_path.stem,
            metadata={
                "source_modified_at": source_mtime.isoformat(),
                "slide_count": slide_count,
                "slides_projected": len(slides_to_project),
                "notes_count": notes_count,
                "adapter_tags": adapter_tags,
                # Declared unconditionally, including when no tags were
                # emitted: the re-ingest merge strips tags matching these
                # prefixes before applying fresh ones, so a run that emits
                # nothing is exactly the run that must clear a stale tag.
                "adapter_tag_prefixes": list(_ADAPTER_TAG_PREFIXES),
            },
        )
